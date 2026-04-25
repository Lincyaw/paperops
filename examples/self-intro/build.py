from __future__ import annotations

import shutil
import sys
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parents[2]
SRC_DIR = REPO_ROOT / "src"
if str(SRC_DIR) not in sys.path:
    sys.path.insert(0, str(SRC_DIR))

from pptx import Presentation as PptxPresentation

from paperops.slides.dsl import Deck, Grid, Heading, KPI, Slide, Subtitle, Text, Title
from paperops.slides.preview import render_slide_preview_powerpoint

OUTPUT_FILE = Path(__file__).with_name("diagnostic_world_model_ir_first.pptx")
SLIDE_TITLES = [
    "Diagnostic intelligence should be trained around world models",
    "A useful diagnostic agent carries state across perception, explanation, and action",
    "Benchmark design should preserve intervention paths",
    "Simulation turns RCA into a controlled learning problem",
    "Trusted diagnosis needs evidence, not just labels",
]
IMPLEMENTED_SLIDE_COUNT = len(SLIDE_TITLES)


def make_theme() -> str:
    return "minimal"


def build_deck() -> Deck:
    deck = Deck(
        theme=make_theme(),
        sheet="seminar",
        meta={"title": "Self intro", "author": "paperops examples", "lang": "en"},
        styles={
            ".summary": {"bg": "bg_alt", "padding": "md", "radius": "md", "border": "border"},
        },
    )
    deck += Slide(class_="cover")[
        Title(SLIDE_TITLES[0]),
        Subtitle("A compact self-intro deck rewritten with the new IR-first API"),
        Text("Trusted diagnosis emerges when the system models causes, observations, and interventions together.", class_="summary"),
    ]
    deck += Slide()[
        Heading(SLIDE_TITLES[1]),
        Grid(style={"cols": "1fr 1fr 1fr", "gap": "md"})[
            Text("Perception: align metrics, logs, traces, and operator context.", class_="summary"),
            Text("Explanation: infer a causal graph instead of outputting a bare label.", class_="summary"),
            Text("Action: recommend the next probe, rollback, or repair with confidence.", class_="summary"),
        ],
    ]
    deck += Slide()[
        Heading(SLIDE_TITLES[2]),
        Grid(style={"cols": "1.4fr 1fr", "gap": "lg"})[
            Text(
                "Benchmarks should store the intervention, the affected services, the observable symptoms, and the expected forward propagation path.",
                class_="summary",
            ),
            Grid(style={"cols": "1fr", "gap": "sm"})[
                KPI(label="Cases", value="9,152", delta="validated", trend="positive"),
                KPI(label="Signals", value="3", delta="metrics/logs/traces", trend="neutral"),
            ],
        ],
    ]
    deck += Slide()[
        Heading(SLIDE_TITLES[3]),
        Text(
            "Simulation gives replayable incidents, controlled interventions, and dense supervision signals for learning why failures unfold.",
            class_="summary",
        ),
    ]
    deck += Slide()[
        Heading(SLIDE_TITLES[4]),
        Grid(style={"cols": "1fr 1fr", "gap": "lg"})[
            Text("Answer quality: find the right failure source.", class_="summary"),
            Text("Reasoning quality: cite evidence that matches the propagation path.", class_="summary"),
            Text("Operational quality: recommend a next action an operator can trust.", class_="summary"),
            Text("Authoring quality: keep the deck in JSON/MDX/Python parity through one shared IR.", class_="summary"),
        ],
    ]
    return deck


def build_presentation(*, output_path: Path | None = None, render_preview: bool = False) -> PptxPresentation:
    destination = output_path or OUTPUT_FILE
    destination.parent.mkdir(parents=True, exist_ok=True)
    build_deck().render(destination)

    if render_preview and shutil.which("soffice") and shutil.which("pdftoppm"):
        preview_dir = destination.with_name("preview")
        preview_dir.mkdir(parents=True, exist_ok=True)
        for stale in preview_dir.glob("slide_*.png"):
            stale.unlink()
        render_slide_preview_powerpoint(str(destination), str(preview_dir))

    return PptxPresentation(str(destination))


def main() -> None:
    build_presentation()


if __name__ == "__main__":
    main()
