"""Slide preview rendering, checking, and integrated deck review."""

from __future__ import annotations

import os
import shutil
import subprocess
import tempfile
from collections import defaultdict
from pathlib import Path

from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt

from paperops.slides.layout.autofit import TextStyle, measure_text_metrics


# ──────────────────────────────────────────────────────────────────────
# Preview rendering (LibreOffice + poppler)
# ──────────────────────────────────────────────────────────────────────

def _find_soffice() -> str:
    """Locate the LibreOffice soffice binary."""
    path = shutil.which("soffice") or shutil.which("libreoffice")
    if path:
        return path
    # macOS default location
    mac_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if os.path.isfile(mac_path):
        return mac_path
    raise RuntimeError(
        "LibreOffice (soffice) not found. Install it:\n"
        "  macOS:  brew install --cask libreoffice\n"
        "  Linux:  sudo apt install libreoffice"
    )


def _find_pdftoppm() -> str:
    """Locate the pdftoppm binary (from poppler)."""
    path = shutil.which("pdftoppm")
    if path:
        return path
    raise RuntimeError(
        "pdftoppm (poppler) not found. Install it:\n"
        "  macOS:  brew install poppler\n"
        "  Linux:  sudo apt install poppler-utils"
    )


def render_slide_preview_powerpoint(
    pptx_path: str, output_dir: str, dpi: int = 200,
) -> list[str]:
    """Render all slides from a pptx file to PNG images.

    Pipeline: PPTX → PDF (via LibreOffice) → PNG per page (via pdftoppm).

    Args:
        pptx_path: Path to the .pptx file.
        output_dir: Directory where PNG files will be saved.
        dpi: Resolution for the exported images (default 200).

    Returns:
        List of paths to the generated PNG files, ordered by slide number.
    """
    soffice = _find_soffice()
    pdftoppm = _find_pdftoppm()

    abs_pptx = os.path.abspath(pptx_path)
    output_dir_path = Path(output_dir)
    output_dir_path.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory() as tmpdir:
        # Step 1: PPTX → PDF
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf",
             "--outdir", tmpdir, abs_pptx],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")

        pdf_name = Path(abs_pptx).stem + ".pdf"
        pdf_path = Path(tmpdir) / pdf_name
        if not pdf_path.is_file():
            raise RuntimeError(f"PDF not created: {pdf_path}")

        # Step 2: PDF → PNG (one per page)
        prefix = Path(tmpdir) / "slide"
        result = subprocess.run(
            [pdftoppm, "-png", "-r", str(dpi), str(pdf_path), str(prefix)],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode != 0:
            raise RuntimeError(f"pdftoppm failed: {result.stderr}")

        # Step 3: collect and rename into output_dir
        png_files = sorted(Path(tmpdir).glob("slide-*.png"))
        output_paths = []
        for idx, png_file in enumerate(png_files, 1):
            dst = output_dir_path / f"slide_{idx:03d}.png"
            shutil.move(str(png_file), str(dst))
            output_paths.append(str(dst))

        return output_paths




# ──────────────────────────────────────────────────────────────────────
# Presentation checker
# ──────────────────────────────────────────────────────────────────────

EMU_PER_INCH = 914400
EMU_PER_PT = 12700


def _estimate_text_height(
    text: str,
    font_size_emu: int,
    box_width_emu: int,
    font_family: str = "Calibri",
    line_spacing: float = 1.0,
) -> int:
    """Estimate text height in EMU using the same text metrics as layout."""
    font_size_pt = font_size_emu / EMU_PER_PT
    box_width_inches = box_width_emu / EMU_PER_INCH if box_width_emu > 0 else None
    metrics = measure_text_metrics(
        text,
        TextStyle(
            font_family=font_family,
            font_size_pt=font_size_pt,
            line_spacing=max(line_spacing, 1.0),
        ),
        max_width_inches=box_width_inches,
    )
    return int(metrics.height * EMU_PER_INCH)


def _overflow_thresholds(
    *,
    text: str,
    usable_width_emu: int,
    line_spacing: float,
    estimated_line_count: int,
) -> tuple[float, int]:
    """Return ratio/absolute thresholds tuned for false-positive reduction."""
    ratio_threshold = 1.12
    absolute_threshold = Inches(0.06)
    is_cjk = _is_cjk_heavy(text)

    if not is_cjk:
        stripped_len = len(text.strip())
        if estimated_line_count <= 1 and stripped_len <= 16:
            ratio_threshold = max(ratio_threshold, 1.35)
            absolute_threshold = max(absolute_threshold, Inches(0.16))

        if usable_width_emu >= Inches(2.8):
            ratio_threshold = 1.18
            absolute_threshold = Inches(0.10)
        elif usable_width_emu >= Inches(1.6):
            ratio_threshold = 1.16
            absolute_threshold = Inches(0.08)

        if estimated_line_count <= 2 and len(text.strip()) <= 64:
            ratio_threshold = max(ratio_threshold, 1.22)
            absolute_threshold = max(absolute_threshold, Inches(0.12))

        if line_spacing <= 1.01 and estimated_line_count <= 2 and len(text.strip()) <= 64:
            ratio_threshold += 0.02
            absolute_threshold = max(absolute_threshold, Inches(0.09))
    else:
        if usable_width_emu <= Inches(1.2):
            ratio_threshold = 1.08
            absolute_threshold = Inches(0.04)

    return ratio_threshold, absolute_threshold


def _get_all_text(shape) -> str:
    parts = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                parts.append(run.text)
            if para.runs:
                parts.append('\n')
    return ''.join(parts).strip()


def _get_font_size(shape) -> int:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if para.font.size:
                return para.font.size
            for run in para.runs:
                if run.font.size:
                    return run.font.size
    return Pt(16)


def _get_font_family(shape) -> str:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if para.font.name:
                return para.font.name
            for run in para.runs:
                if run.font.name:
                    return run.font.name
    return "Calibri"


def _get_line_spacing(shape) -> float:
    if not shape.has_text_frame:
        return 1.0

    default_font_size = _get_font_size(shape)
    default_pt = default_font_size.pt if hasattr(default_font_size, "pt") else float(default_font_size)
    for para in shape.text_frame.paragraphs:
        value = para.line_spacing
        if value is None:
            continue
        if hasattr(value, "pt"):
            if default_pt > 0:
                return max(value.pt / default_pt, 0.8)
            continue
        try:
            numeric = float(value)
        except (TypeError, ValueError):
            continue
        if 0.1 <= numeric <= 4.0:
            return max(numeric, 0.8)
    return 1.0


def _get_shape_bounds(shape):
    left = shape.left or 0
    top = shape.top or 0
    width = shape.width or 0
    height = shape.height or 0
    return (left, top, left + width, top + height)


def _shapes_overlap(b1, b2, tolerance_emu=0):
    l1, t1, r1, bot1 = b1
    l2, t2, r2, bot2 = b2
    l1 += tolerance_emu
    t1 += tolerance_emu
    r1 -= tolerance_emu
    bot1 -= tolerance_emu
    overlap_left = max(l1, l2)
    overlap_top = max(t1, t2)
    overlap_right = min(r1, r2)
    overlap_bottom = min(bot1, bot2)
    if overlap_left < overlap_right and overlap_top < overlap_bottom:
        return (overlap_right - overlap_left) * (overlap_bottom - overlap_top)
    return 0


def _emu_to_inches(emu):
    return emu / EMU_PER_INCH


def _bounds_to_inches(bounds):
    return tuple(_emu_to_inches(v) for v in bounds)


def _shape_info(shape) -> str:
    text = _get_all_text(shape)
    if text:
        preview = text[:50].replace('\n', ' ')
        if len(text) > 50:
            preview += '...'
        return f'"{preview}"'
    return f'<shape>'


def _longest_token_len(text: str) -> int:
    tokens = [token for token in text.replace("\n", " ").split(" ") if token]
    if not tokens:
        return len(text.strip())
    return max(len(token) for token in tokens)


def _measure_preview_text_width(draw, text: str, font) -> int:
    if not text:
        return 0
    left, _top, right, _bottom = draw.textbbox((0, 0), text, font=font)
    return max(right - left, 0)


def _is_cjk_heavy(text: str) -> bool:
    if not text:
        return False
    total = 0
    cjk = 0
    for ch in text:
        if ch.isspace():
            continue
        total += 1
        cp = ord(ch)
        if (
            0x4E00 <= cp <= 0x9FFF
            or 0x3400 <= cp <= 0x4DBF
            or 0xF900 <= cp <= 0xFAFF
            or 0x3000 <= cp <= 0x303F
        ):
            cjk += 1
    if total == 0:
        return False
    return cjk / total >= 0.4


def _wrap_preview_chars(draw, text: str, font, max_width_px: int) -> list[str]:
    lines: list[str] = []
    current = ""
    for ch in text:
        candidate = current + ch
        if _measure_preview_text_width(draw, candidate, font) <= max_width_px or not current:
            current = candidate
            continue
        lines.append(current)
        current = ch
    if current:
        lines.append(current)
    return lines or [text]


def _wrap_preview_text(draw, text: str, font, max_width_px: int) -> list[str]:
    """Wrap text for preview rendering and diagnostics.

    Handles both whitespace-delimited text and CJK-heavy text where spaces may
    be sparse or absent.
    """
    if not text:
        return []

    wrapped: list[str] = []
    for paragraph in text.splitlines() or [""]:
        if not paragraph.strip():
            wrapped.append("")
            continue

        if _is_cjk_heavy(paragraph) or " " not in paragraph:
            wrapped.extend(_wrap_preview_chars(draw, paragraph, font, max_width_px))
            continue

        tokens = [t for t in paragraph.split(" ") if t]
        line = ""
        for token in tokens:
            candidate = f"{line} {token}".strip() if line else token
            if _measure_preview_text_width(draw, candidate, font) <= max_width_px:
                line = candidate
                continue

            if line:
                wrapped.append(line)
                line = ""

            if _measure_preview_text_width(draw, token, font) > max_width_px:
                token_lines = _wrap_preview_chars(draw, token, font, max_width_px)
                wrapped.extend(token_lines[:-1])
                line = token_lines[-1]
            else:
                line = token

        if line:
            wrapped.append(line)

    return wrapped


def summarize_slide_shapes(slide) -> dict:
    """Return lightweight, deterministic heuristics about slide density."""
    summary = {
        "shape_count": len(slide.shapes),
        "text_shape_count": 0,
        "small_text_shape_count": 0,
        "narrow_text_shapes": [],
        "crowding_risks": [],
        "density": "normal",
    }

    for idx, shape in enumerate(slide.shapes, 1):
        try:
            has_text = shape.has_text_frame
        except Exception:
            has_text = False
        if not has_text:
            continue

        text = _get_all_text(shape)
        if not text:
            continue

        summary["text_shape_count"] += 1
        width_inches = _emu_to_inches(shape.width or 0)
        height_inches = _emu_to_inches(shape.height or 0)
        longest = _longest_token_len(text)
        ratio = width_inches / max(longest, 1)
        shape_info = {
            "shape_index": idx,
            "text_preview": text[:60],
            "width_inches": round(width_inches, 3),
            "height_inches": round(height_inches, 3),
            "token_width_ratio": round(ratio, 3),
        }

        if width_inches < 1.2 or height_inches < 0.45:
            summary["small_text_shape_count"] += 1
        if ratio < 0.09:
            summary["narrow_text_shapes"].append(shape_info)
        if ratio < 0.07 or (width_inches < 1.5 and longest >= 12):
            summary["crowding_risks"].append(shape_info)

    if summary["text_shape_count"] >= 8 or summary["small_text_shape_count"] >= 4:
        summary["density"] = "high"
    elif summary["text_shape_count"] <= 3 and not summary["crowding_risks"]:
        summary["density"] = "low"

    return summary


def _normalize_issue(
    issue: dict,
    *,
    source: str,
    default_slide: int | None = None,
    default_severity: str = "warning",
) -> dict:
    item = dict(issue)
    if item.get("slide") is None and default_slide is not None:
        item["slide"] = default_slide
    if "source" not in item:
        item["source"] = source
    if "code" not in item:
        item["code"] = item.get("type", f"{source}_issue")
    if "type" not in item:
        item["type"] = item["code"]
    if "message" not in item:
        item["message"] = item.get("detail", item["code"])
    if "detail" not in item:
        item["detail"] = item["message"]
    if "severity" not in item:
        item["severity"] = default_severity
    return item


def summarize_presentation_issues(issues: list[dict]) -> dict[int, list[dict]]:
    grouped: dict[int, list[dict]] = defaultdict(list)
    for issue in issues:
        slide_num = issue.get("slide")
        if isinstance(slide_num, int):
            grouped[slide_num].append(issue)
    return dict(grouped)


def _count_issues(issues: list[dict]) -> dict[str, int]:
    counts = {
        "total": 0,
        "layout": 0,
        "saved_file": 0,
        "style": 0,
        "error": 0,
        "warning": 0,
        "info": 0,
    }
    for issue in issues:
        counts["total"] += 1
        source = issue.get("source")
        if source in counts:
            counts[source] += 1
        severity = issue.get("severity")
        if severity in counts:
            counts[severity] += 1
    return counts


def _score_slide_summary(summary: dict, issues: list[dict]) -> int:
    layout_issues = [i for i in issues if i.get("source") == "layout"]
    saved_issues = [i for i in issues if i.get("source") == "saved_file"]
    return (
        len(layout_issues) * 4
        + len(saved_issues) * 5
        + len(summary.get("crowding_risks", [])) * 3
        + len(summary.get("narrow_text_shapes", [])) * 2
        + (2 if summary.get("density") == "high" else 0)
    )


def check_presentation(pptx_path: str) -> list[dict]:
    """Open a .pptx file and check for layout issues.

    Returns list of issue dicts with keys: slide, type, detail, etc.
    """
    prs = PptxPresentation(pptx_path)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    all_issues = []

    for slide_num, slide in enumerate(prs.slides, 1):
        shapes_with_text = []

        for shape in slide.shapes:
            bounds = _get_shape_bounds(shape)
            info = _shape_info(shape)

            # Check off-slide
            l, t, r, b = bounds
            if (l < -Inches(0.1) or t < -Inches(0.1) or
                    r > slide_width + Inches(0.1) or b > slide_height + Inches(0.1)):
                all_issues.append({
                    'slide': slide_num,
                    'source': 'saved_file',
                    'severity': 'error',
                    'code': 'off_slide',
                    'type': 'off_slide',
                    'message': f'Shape {info} extends beyond slide boundaries',
                    'detail': f'Shape {info} extends beyond slide boundaries',
                    'position': _bounds_to_inches(bounds),
                })

            # Check text overflow
            if shape.has_text_frame:
                text = _get_all_text(shape)
                if text:
                    font_size = _get_font_size(shape)
                    box_w = shape.width or 0
                    box_h = shape.height or 0

                    tf = shape.text_frame
                    margin_l = tf.margin_left or Inches(0.05)
                    margin_r = tf.margin_right or Inches(0.05)
                    margin_t = tf.margin_top or Inches(0.05)
                    margin_b = tf.margin_bottom or Inches(0.05)

                    usable_w = box_w - margin_l - margin_r
                    usable_h = box_h - margin_t - margin_b

                    if usable_w > 0 and usable_h > 0:
                        line_spacing = _get_line_spacing(shape)
                        is_cjk = _is_cjk_heavy(text)
                        width_slack = Inches(0.06) if is_cjk else Inches(0.12)
                        estimated_metrics = measure_text_metrics(
                            text,
                            TextStyle(
                                font_family=_get_font_family(shape),
                                font_size_pt=(font_size.pt if hasattr(font_size, "pt") else float(font_size)),
                                line_spacing=max(line_spacing, 0.95),
                            ),
                            max_width_inches=(usable_w + width_slack) / EMU_PER_INCH,
                        )
                        needed_h = _estimate_text_height(
                            text,
                            font_size,
                            usable_w + width_slack,
                            font_family=_get_font_family(shape),
                            line_spacing=line_spacing,
                        )
                        ratio_threshold, absolute_threshold = _overflow_thresholds(
                            text=text,
                            usable_width_emu=usable_w,
                            line_spacing=line_spacing,
                            estimated_line_count=max(estimated_metrics.line_count, 1),
                        )
                        overflow_gap = needed_h - usable_h
                        if needed_h > usable_h * ratio_threshold and overflow_gap > absolute_threshold:
                            overflow_pct = (needed_h / usable_h - 1) * 100
                            all_issues.append({
                                'slide': slide_num,
                                'source': 'saved_file',
                                'severity': 'warning',
                                'code': 'text_overflow',
                                'type': 'text_overflow',
                                'message': f'Text {info} overflows by {overflow_pct:.0f}%',
                                'detail': (f'Text {info} overflows by '
                                           f'{overflow_pct:.0f}%'),
                                'position': _bounds_to_inches(bounds),
                            })

                    shapes_with_text.append((shape, bounds, info))

        # Check overlaps between text shapes
        for i in range(len(shapes_with_text)):
            for j in range(i + 1, len(shapes_with_text)):
                s1, b1, info1 = shapes_with_text[i]
                s2, b2, info2 = shapes_with_text[j]
                overlap = _shapes_overlap(b1, b2, tolerance_emu=Inches(0.05))
                if overlap > 0:
                    area1 = (b1[2] - b1[0]) * (b1[3] - b1[1])
                    area2 = (b2[2] - b2[0]) * (b2[3] - b2[1])
                    smaller = min(area1, area2)
                    if smaller > 0:
                        pct = overlap / smaller * 100
                        if pct > 10:
                            all_issues.append({
                                'slide': slide_num,
                                'source': 'saved_file',
                                'severity': 'warning',
                                'code': 'overlap',
                                'type': 'overlap',
                                'message': f'{info1} and {info2} overlap by {pct:.0f}%',
                                'detail': (f'{info1} and {info2} overlap by '
                                           f'{pct:.0f}%'),
                            })

    return all_issues


def review_deck_artifacts(
    pptx_path: str,
    layout_issues: list[dict] | None = None,
    preview_paths: list[str] | None = None,
    slide_titles: list[str] | None = None,
    layout_summary: dict | None = None,
) -> dict:
    """Return a stable, merged review artifact for a generated deck."""
    prs = PptxPresentation(pptx_path)
    raw_layout_issues = layout_issues or []
    layout_issues = [
        _normalize_issue(issue, source="layout")
        for issue in raw_layout_issues
        if isinstance(issue, dict)
    ]
    saved_issues = [
        _normalize_issue(issue, source="saved_file")
        for issue in check_presentation(pptx_path)
    ]
    preview_paths = preview_paths or []
    slide_titles = slide_titles or []
    layout_summary = layout_summary or {}

    layout_by_slide = summarize_presentation_issues(layout_issues)
    saved_by_slide = summarize_presentation_issues(saved_issues)
    preview_by_slide = {
        idx + 1: path
        for idx, path in enumerate(preview_paths)
    }

    slides: list[dict] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        summary = summarize_slide_shapes(slide)
        slide_layout_issues = layout_by_slide.get(slide_num, [])
        slide_saved_issues = saved_by_slide.get(slide_num, [])
        combined_issues = [*slide_layout_issues, *slide_saved_issues]
        slide_entry = {
            "slide_number": slide_num,
            "title": slide_titles[slide_num - 1] if slide_num - 1 < len(slide_titles) else None,
            "issues": combined_issues,
            "layout_issues": slide_layout_issues,
            "saved_file_issues": slide_saved_issues,
            "preview_path": preview_by_slide.get(slide_num),
            "summary": summary,
        }
        slide_entry["issue_count"] = len(combined_issues)
        slide_entry["issue_counts"] = _count_issues(combined_issues)
        slide_entry["score"] = _score_slide_summary(summary, combined_issues)
        slides.append(slide_entry)

    ranked = sorted(
        (
            {
                "slide_number": slide["slide_number"],
                "title": slide["title"],
                "score": slide["score"],
                "issue_count": slide["issue_count"],
            }
            for slide in slides
            if slide["issue_count"] > 0
        ),
        key=lambda item: (-item["score"], item["slide_number"]),
    )

    all_issues = [*layout_issues, *saved_issues]
    issue_counts = _count_issues(all_issues)

    return {
        "schema_version": "2026-04-09",
        "artifact": "deck_review",
        "total_slides": len(prs.slides),
        "issue_counts": issue_counts,
        "issues": all_issues,
        "preview_paths": preview_paths,
        "top_problem_slides": ranked,
        "slides": slides,
        "layout_summary": layout_summary,
        # Legacy compatibility fields.
        "layout_issue_count": len(layout_issues),
        "saved_issue_count": len(saved_issues),
    }
