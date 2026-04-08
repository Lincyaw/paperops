"""SlideBuilder — renders a component tree onto a python-pptx slide."""

from __future__ import annotations

import os
import sys
import tempfile

from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

from paperops.slides.core.constants import (
    Region, SLIDE_WIDTH, SLIDE_HEIGHT,
    CONTENT_REGION, TITLE_REGION, REFERENCE_REGION, ACCENT_LINE_TOP,
    Direction, Align,
)
from paperops.slides.core.types import resolve_color as _resolve_color_raw
from paperops.slides.core.types import resolve_font_size, resolve_size
from paperops.slides.layout.engine import compute_layout
from paperops.slides.layout.containers import LayoutNode, HStack, VStack, Grid, Padding
from paperops.slides.components.shapes import Box, RoundedBox, Circle, Badge, Arrow, Line
from paperops.slides.components.text import TextBlock, BulletList
from paperops.slides.components.table import Table
from paperops.slides.components.image import Image, SvgImage
from paperops.slides.components.composite import Callout, Flow
from paperops.slides.components.charts.bar import BarChart
from paperops.slides.components.charts.radar import RadarChart
from paperops.slides.components.charts.flowchart import Flowchart
from paperops.slides.components.charts.line import LineChart
from paperops.slides.components.charts.pie import PieChart
from paperops.slides.components.charts.horizontal_bar import HorizontalBarChart
from paperops.slides.animation import inject_appear_animations

# Color literals that components use but aren't in themes
_LITERAL_COLORS = {
    "white": "#FFFFFF",
    "black": "#000000",
}


def resolve_color(theme, color):
    """Wrap resolve_color to handle literal color names like 'white'/'black'."""
    if isinstance(color, str) and color in _LITERAL_COLORS:
        h = _LITERAL_COLORS[color].lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    return _resolve_color_raw(theme, color)


_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


class SlideBuilder:
    """Builder for a single slide. Created by Presentation.slide()."""

    def __init__(self, pptx_slide, theme, title=None, reference=None):
        self._slide = pptx_slide
        self._theme = theme
        self._title = title
        self._reference = reference
        self._component = None
        self._click_groups = None
        self._component_to_shape_ids: dict[int, list[int]] = {}  # id(component) -> [shape_ids]
        self._issues: list[dict] = []
        self._rendered = False
        self._bg_color = None
        self._bg_image = None

    def layout(self, component):
        """Set the content layout. Component can be any LayoutNode or container."""
        self._component = component
        return self

    def notes(self, text: str):
        """Set speaker notes for this slide."""
        notes_slide = self._slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = text
        return self

    def animate(self, click_groups: list[list]):
        """Declare click-to-advance animation groups.

        click_groups: list of lists of component objects.
        """
        self._click_groups = click_groups
        return self

    def background(self, color=None, image_path=None):
        """Set slide background color or image."""
        self._bg_color = color
        self._bg_image = image_path
        return self

    # ------------------------------------------------------------------
    # Rendering
    # ------------------------------------------------------------------

    def _render(self):
        """Render everything onto the pptx slide."""
        if self._rendered:
            return self._issues
        self._rendered = True

        # 0. Background
        if self._bg_image:
            self._slide.shapes.add_picture(
                self._bg_image,
                Inches(0), Inches(0),
                Inches(SLIDE_WIDTH), Inches(SLIDE_HEIGHT),
            )
        elif self._bg_color:
            bg_fill = self._slide.background.fill
            bg_fill.solid()
            bg_fill.fore_color.rgb = resolve_color(self._theme, self._bg_color)

        # 1. Title
        if self._title:
            self._render_title()

        # 2. Reference
        if self._reference:
            self._render_reference()

        # 3. Layout engine + render component tree
        if self._component is not None:
            issues = compute_layout(self._component, CONTENT_REGION, self._theme)
            self._issues.extend(issues)
            self._render_node(self._component)

        # 4. Animations
        if self._click_groups:
            resolved = []
            for group in self._click_groups:
                ids = []
                for comp in group:
                    # Collect all leaf shape IDs for this component,
                    # whether it's a leaf, container, or composite
                    ids.extend(self._collect_shape_ids(comp))
                if ids:
                    resolved.append(ids)
            if resolved:
                # Shapes in groups 2+ need to start hidden
                initially_hidden = []
                for group_ids in resolved[1:]:
                    initially_hidden.extend(group_ids)
                inject_appear_animations(
                    self._slide._element, resolved,
                    initially_hidden=initially_hidden,
                )

        return self._issues

    # ------------------------------------------------------------------
    # Title / Reference helpers
    # ------------------------------------------------------------------

    def _render_title(self):
        r = TITLE_REGION
        txbox = self._slide.shapes.add_textbox(
            Inches(r.left), Inches(r.top), Inches(r.width), Inches(r.height),
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = self._title
        p.font.size = resolve_font_size(self._theme, "title")
        p.font.color.rgb = resolve_color(self._theme, "text")
        p.font.bold = True
        p.font.name = self._theme.font_family

        # Accent line below title
        line_y = ACCENT_LINE_TOP
        line_shape = self._slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(r.left), Inches(line_y),
            Inches(1.2), Inches(0.04),
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = resolve_color(self._theme, "primary")
        line_shape.line.fill.background()

    def _render_reference(self):
        r = REFERENCE_REGION
        txbox = self._slide.shapes.add_textbox(
            Inches(r.left), Inches(r.top), Inches(r.width), Inches(r.height),
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = self._reference
        p.font.size = resolve_font_size(self._theme, "small")
        p.font.color.rgb = resolve_color(self._theme, "text_light")
        p.font.name = self._theme.font_family
        p.alignment = PP_ALIGN.RIGHT

    # ------------------------------------------------------------------
    # Component rendering dispatch
    # ------------------------------------------------------------------

    def _render_node(self, node: LayoutNode):
        """Recursively render a positioned node."""
        # Containers: recurse into children
        if isinstance(node, (HStack, VStack, Grid)):
            for child in node.children:
                self._render_node(child)
            return
        if isinstance(node, Padding):
            if node.child is not None:
                self._render_node(node.child)
            return

        # Composites: expand then render, track child shapes for parent
        if isinstance(node, (Callout, Flow, Flowchart)):
            expanded = node.to_layout()
            # Re-layout the expanded tree within the node's region
            if node._region is not None:
                compute_layout(expanded, node._region, self._theme)
                self._render_node(expanded)
                self._track_composite(node, expanded)
            return

        # Charts: render to SVG then as image
        if isinstance(node, (BarChart, RadarChart, LineChart, PieChart, HorizontalBarChart)):
            svg_str = node.to_svg(self._theme)
            self._render_svg_image(node, svg_str)
            return

        # Leaf components
        if isinstance(node, Table):
            self._render_table(node)
        elif isinstance(node, BulletList):
            self._render_bullet_list(node)
        elif isinstance(node, TextBlock):
            self._render_text_block(node)
        elif isinstance(node, Badge):
            self._render_badge(node)
        elif isinstance(node, Circle):
            self._render_circle(node)
        elif isinstance(node, RoundedBox):
            self._render_rounded_box(node)
        elif isinstance(node, Box):
            self._render_box(node)
        elif isinstance(node, Arrow):
            self._render_arrow(node)
        elif isinstance(node, Line):
            self._render_line(node)
        elif isinstance(node, SvgImage):
            svg_str = node.svg
            if hasattr(svg_str, 'render'):
                svg_str = svg_str.render()
            self._render_svg_image(node, svg_str)
        elif isinstance(node, Image):
            self._render_image(node)

    # ------------------------------------------------------------------
    # Leaf renderers
    # ------------------------------------------------------------------

    def _region_args(self, node):
        """Return (left, top, width, height) in EMU from node._region."""
        r = node._region
        return Inches(r.left), Inches(r.top), Inches(r.width), Inches(r.height)

    def _track(self, node, shape):
        """Record component -> shape_id mapping for animation."""
        spid = shape.shape_id
        self._component_to_shape_ids.setdefault(id(node), []).append(spid)

    def _track_composite(self, composite_node, expanded_node):
        """After rendering a composite, copy all child shape IDs to the composite."""
        child_ids = self._collect_shape_ids(expanded_node)
        if child_ids:
            self._component_to_shape_ids.setdefault(id(composite_node), []).extend(child_ids)

    def _collect_shape_ids(self, node) -> list[int]:
        """Collect all shape IDs from a rendered subtree."""
        ids = []
        cid = id(node)
        if cid in self._component_to_shape_ids:
            ids.extend(self._component_to_shape_ids[cid])
        if isinstance(node, (HStack, VStack, Grid)):
            for child in node.children:
                ids.extend(self._collect_shape_ids(child))
        elif isinstance(node, Padding) and node.child is not None:
            ids.extend(self._collect_shape_ids(node.child))
        return ids

    def _render_shape(self, node, mso_shape, *, has_border=True, region_override=None):
        """Shared renderer for Box, RoundedBox, Circle, Badge."""
        if region_override is not None:
            left, top = Inches(region_override.left), Inches(region_override.top)
            width, height = Inches(region_override.width), Inches(region_override.height)
        else:
            left, top, width, height = self._region_args(node)
        shape = self._slide.shapes.add_shape(mso_shape, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = resolve_color(self._theme, node.color)
        if has_border and hasattr(node, 'border'):
            shape.line.color.rgb = resolve_color(self._theme, node.border)
            shape.line.width = Pt(node.border_width)
        else:
            shape.line.fill.background()
        text = getattr(node, 'text', '')
        if text:
            tf = shape.text_frame
            # Smart word wrap: disable for short text that should fit on one line
            # This prevents unwanted line breaks like "Telemet ry" instead of "Telemetry"
            font_size_val = resolve_font_size(self._theme, node.font_size)
            if font_size_val is None:
                font_size_pt = 14.0  # Default to 14pt (caption size)
            elif hasattr(font_size_val, 'pt'):
                font_size_pt = font_size_val.pt
            else:
                font_size_pt = float(font_size_val)
            # Estimate text width: char width is roughly font_size_pt / 72 * 0.6 for typical fonts
            char_width_inches = (font_size_pt / 72) * 0.6
            estimated_text_width = len(text) * char_width_inches
            # Get box width from the shape we just created
            box_width = width.inches
            # Add safety margin for text frame margins (typically 0.1-0.2 inches)
            usable_width = box_width - 0.25
            # Disable word wrap if text should fit in one line
            # Use 1.0x factor since we already accounted for margins in usable_width
            should_wrap = estimated_text_width > usable_width
            tf.word_wrap = should_wrap
            align = _ALIGN_MAP.get(getattr(node, 'align', 'center'), PP_ALIGN.CENTER)
            p = tf.paragraphs[0]
            p.text = text
            p.alignment = align
            p.font.size = Pt(font_size_pt)
            p.font.color.rgb = resolve_color(self._theme, node.text_color)
            p.font.bold = node.bold
            p.font.name = self._theme.font_family
        self._track(node, shape)

    def _render_box(self, node: Box):
        self._render_shape(node, MSO_SHAPE.RECTANGLE)

    def _render_rounded_box(self, node: RoundedBox):
        self._render_shape(node, MSO_SHAPE.ROUNDED_RECTANGLE)

    def _render_circle(self, node: Circle):
        # Force circle to be square: use min dimension, center within region
        r = node._region
        square_region = None
        if r and abs(r.width - r.height) > 0.01:
            d = min(r.width, r.height)
            cx = r.left + (r.width - d) / 2
            cy = r.top + (r.height - d) / 2
            square_region = Region(left=cx, top=cy, width=d, height=d)
        self._render_shape(node, MSO_SHAPE.OVAL, has_border=False,
                           region_override=square_region)

    def _render_badge(self, node: Badge):
        self._render_shape(node, MSO_SHAPE.ROUNDED_RECTANGLE, has_border=False)

    def _render_text_block(self, node: TextBlock):
        left, top, width, height = self._region_args(node)
        txbox = self._slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = True

        alignment = _ALIGN_MAP.get(node.align, PP_ALIGN.LEFT)
        default_font_size = resolve_font_size(self._theme, node.font_size)
        default_color = resolve_color(self._theme, node.color)

        if node.runs is not None:
            # RichText mode: render each (text, format_dict) as a separate run
            p = tf.paragraphs[0]
            p.alignment = alignment
            for run_text, fmt in node.runs:
                run = p.add_run()
                run.text = run_text
                run.font.name = self._theme.font_family
                run.font.size = resolve_font_size(self._theme, fmt.get('font_size', node.font_size))
                run.font.bold = fmt.get('bold', node.bold)
                run.font.italic = fmt.get('italic', node.italic)
                if 'color' in fmt:
                    run.font.color.rgb = resolve_color(self._theme, fmt['color'])
                else:
                    run.font.color.rgb = default_color
        else:
            # Plain text mode: split on newlines to create proper paragraphs
            lines = node.text.split('\n')
            for i, line in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.alignment = alignment
                p.font.size = default_font_size
                p.font.color.rgb = default_color
                p.font.bold = node.bold
                p.font.italic = node.italic
                p.font.name = self._theme.font_family

        self._track(node, txbox)

    def _render_bullet_list(self, node: BulletList):
        left, top, width, height = self._region_args(node)
        txbox = self._slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(node.items):
            if isinstance(item, tuple):
                text, indent = item
            else:
                text = item
                indent = 0

            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = text
            p.level = indent
            p.font.size = resolve_font_size(self._theme, node.font_size)
            p.font.color.rgb = resolve_color(self._theme, node.color)
            p.font.name = self._theme.font_family
            # Bullet character
            pPr = p._pPr
            if pPr is None:
                from lxml import etree
                nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                pPr = etree.SubElement(
                    p._p,
                    '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr',
                )
            buChar = pPr.makeelement(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar', {}
            )
            buChar.set('char', '\u2022')
            pPr.append(buChar)

        self._track(node, txbox)

    def _render_table(self, node: Table):
        r = node._region
        num_rows = len(node.rows) + (1 if node.headers else 0)
        num_cols = len(node.headers) if node.headers else (len(node.rows[0]) if node.rows else 1)

        table_shape = self._slide.shapes.add_table(
            num_rows, num_cols,
            Inches(r.left), Inches(r.top), Inches(r.width), Inches(r.height),
        )
        tbl = table_shape.table

        font_size = resolve_font_size(self._theme, node.font_size)
        header_fill = resolve_color(self._theme, node.header_color)
        header_text = resolve_color(self._theme, node.header_text_color)
        text_color = resolve_color(self._theme, "text")
        bg_alt = resolve_color(self._theme, "bg_alt")
        bg_color = resolve_color(self._theme, "bg")

        row_offset = 0
        if node.headers:
            for ci, hdr in enumerate(node.headers):
                cell = tbl.cell(0, ci)
                cell.text = hdr
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                for p in cell.text_frame.paragraphs:
                    p.font.size = font_size
                    p.font.color.rgb = header_text
                    p.font.bold = True
                    p.font.name = self._theme.font_family
                    p.alignment = PP_ALIGN.CENTER
            row_offset = 1

        for ri, row_data in enumerate(node.rows):
            for ci, cell_text in enumerate(row_data):
                cell = tbl.cell(ri + row_offset, ci)
                cell.text = str(cell_text)
                # Alternating row colors
                fill_rgb = bg_alt if ri % 2 == 0 else bg_color
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill_rgb
                for p in cell.text_frame.paragraphs:
                    p.font.size = font_size
                    p.font.color.rgb = text_color
                    p.font.name = self._theme.font_family

        self._track(node, table_shape)

    def _resolve_connector_endpoints(self, node):
        """Resolve from/to regions for Arrow/Line. Returns (x1, y1, x2, y2) or None."""
        fr = node.from_component._region if node.from_component and node.from_component._region else None
        tr = node.to_component._region if node.to_component and node.to_component._region else None
        direction = getattr(node, 'direction', 'horizontal')
        if fr and tr:
            if direction == "vertical":
                # Connect bottom-center of source to top-center of target
                return (fr.left + fr.width / 2, fr.bottom,
                        tr.left + tr.width / 2, tr.top)
            else:
                # Connect right-center of source to left-center of target
                return fr.right, fr.top + fr.height / 2, tr.left, tr.top + tr.height / 2
        if node._region:
            r = node._region
            if direction == "vertical":
                return r.left + r.width / 2, r.top, r.left + r.width / 2, r.bottom
            return r.left, r.top + r.height / 2, r.right, r.top + r.height / 2
        return None

    def _render_arrow(self, node: Arrow):
        pts = self._resolve_connector_endpoints(node)
        if not pts:
            return
        x1, y1, x2, y2 = pts
        connector = self._slide.shapes.add_connector(
            1, Inches(x1), Inches(y1), Inches(x2), Inches(y2),
        )
        connector.line.color.rgb = resolve_color(self._theme, node.color)
        connector.line.width = Pt(node.width_pt)
        from pptx.oxml.ns import qn
        ln = connector.line._ln
        tailEnd = ln.makeelement(qn('a:tailEnd'), {})
        tailEnd.set('type', 'triangle')
        tailEnd.set('w', 'med')
        tailEnd.set('len', 'med')
        ln.append(tailEnd)
        self._track(node, connector)

        # Render arrow label at midpoint if provided
        if node.label is not None:
            mid_x = (x1 + x2) / 2
            mid_y = (y1 + y2) / 2
            label_w = 1.5
            label_h = 0.3
            direction = getattr(node, 'direction', 'horizontal')
            if direction == 'vertical':
                # Offset to the right for vertical arrows
                lbl_left = mid_x + 0.1
                lbl_top = mid_y - label_h / 2
            else:
                # Offset above for horizontal arrows
                lbl_left = mid_x - label_w / 2
                lbl_top = mid_y - label_h - 0.1
            lbl_box = self._slide.shapes.add_textbox(
                Inches(lbl_left), Inches(lbl_top),
                Inches(label_w), Inches(label_h),
            )
            lbl_tf = lbl_box.text_frame
            lbl_tf.word_wrap = True
            lbl_p = lbl_tf.paragraphs[0]
            lbl_p.text = node.label
            lbl_p.alignment = PP_ALIGN.CENTER
            lbl_p.font.size = resolve_font_size(self._theme, "caption")
            lbl_p.font.color.rgb = resolve_color(self._theme, "text_mid")
            lbl_p.font.name = self._theme.font_family
            self._track(node, lbl_box)

    def _render_line(self, node: Line):
        pts = self._resolve_connector_endpoints(node)
        if not pts:
            return
        x1, y1, x2, y2 = pts
        connector = self._slide.shapes.add_connector(
            1, Inches(x1), Inches(y1), Inches(x2), Inches(y2),
        )
        connector.line.color.rgb = resolve_color(self._theme, node.color)
        connector.line.width = Pt(node.width_pt)
        if node.dashed:
            from pptx.oxml.ns import qn
            ln = connector.line._ln
            prstDash = ln.makeelement(qn('a:prstDash'), {})
            prstDash.set('val', 'dash')
            ln.append(prstDash)
        self._track(node, connector)

    def _render_image(self, node: Image):
        from PIL import Image as PILImage

        r = node._region
        if not getattr(node, "preserve_aspect", True):
            pic = self._slide.shapes.add_picture(
                node.path,
                Inches(r.left), Inches(r.top), Inches(r.width), Inches(r.height),
            )
            self._track(node, pic)
            return

        with PILImage.open(node.path) as img:
            img_w, img_h = img.size

        img_aspect = img_w / img_h if img_h > 0 else 1.0
        region_aspect = r.width / r.height if r.height > 0 else 1.0

        if img_aspect > region_aspect:
            fit_w = r.width
            fit_h = r.width / img_aspect
        else:
            fit_h = r.height
            fit_w = r.height * img_aspect

        fit_left = r.left + (r.width - fit_w) / 2
        fit_top = r.top + (r.height - fit_h) / 2

        pic = self._slide.shapes.add_picture(
            node.path,
            Inches(fit_left), Inches(fit_top), Inches(fit_w), Inches(fit_h),
        )
        self._track(node, pic)

    def _render_svg_image(self, node, svg_str: str):
        """Render SVG string to PNG and add as picture, preserving aspect ratio."""
        import cairosvg
        import struct

        r = node._region
        scale = getattr(node, 'scale', 3)

        png_bytes = cairosvg.svg2png(bytestring=svg_str.encode('utf-8'), scale=scale)

        # Read dimensions from PNG header (bytes 16-24) instead of decoding
        img_w, img_h = struct.unpack('>II', png_bytes[16:24])
        img_aspect = img_w / img_h if img_h > 0 else 1.0

        region_aspect = r.width / r.height if r.height > 0 else 1.0

        # Fit image within region, preserving aspect ratio, centered
        if img_aspect > region_aspect:
            # Image is wider than region — fit to width, shrink height
            fit_w = r.width
            fit_h = r.width / img_aspect
        else:
            # Image is taller than region — fit to height, shrink width
            fit_h = r.height
            fit_w = r.height * img_aspect

        # Center within the allocated region
        fit_left = r.left + (r.width - fit_w) / 2
        fit_top = r.top + (r.height - fit_h) / 2

        tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
        tmp.write(png_bytes)
        tmp.close()
        try:
            pic = self._slide.shapes.add_picture(
                tmp.name,
                Inches(fit_left), Inches(fit_top), Inches(fit_w), Inches(fit_h),
            )
            self._track(node, pic)
        finally:
            os.unlink(tmp.name)
