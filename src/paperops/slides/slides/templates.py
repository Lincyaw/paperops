"""High-level slide templates — registered onto the Presentation class."""

from __future__ import annotations

from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from paperops.slides.core.constants import (
    Region, SLIDE_WIDTH, SLIDE_HEIGHT,
    CONTENT_REGION, TITLE_REGION, ACCENT_LINE_TOP,
    MARGIN_LEFT, MARGIN_RIGHT,
)
from paperops.slides.core.types import resolve_color, resolve_font_size
from paperops.slides.layout.containers import HStack, VStack, Padding
from paperops.slides.components.shapes import Box, RoundedBox
from paperops.slides.components.text import TextBlock, BulletList
from paperops.slides.components.table import Table


def register_templates(presentation_class):
    """Add high-level template methods to Presentation."""

    def cover(self, title, subtitle="", author=""):
        """Cover slide — centered title, accent line, subtitle."""
        sb = self.slide()
        slide = sb._slide

        # Center region for title
        cx = SLIDE_WIDTH / 2
        title_w = SLIDE_WIDTH - 2 * MARGIN_LEFT

        # Title text
        txbox = slide.shapes.add_textbox(
            Inches(MARGIN_LEFT), Inches(2.0),
            Inches(title_w), Inches(1.5),
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(40)
        p.font.color.rgb = resolve_color(self._theme, "text")
        p.font.bold = True
        p.font.name = self._theme.font_family

        # Accent line
        line_w = 2.0
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(cx - line_w / 2), Inches(3.7),
            Inches(line_w), Inches(0.05),
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = resolve_color(self._theme, "primary")
        line_shape.line.fill.background()

        # Subtitle
        if subtitle:
            txbox2 = slide.shapes.add_textbox(
                Inches(MARGIN_LEFT), Inches(4.0),
                Inches(title_w), Inches(1.0),
            )
            tf2 = txbox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.alignment = PP_ALIGN.CENTER
            p2.font.size = resolve_font_size(self._theme, "subtitle")
            p2.font.color.rgb = resolve_color(self._theme, "text_mid")
            p2.font.name = self._theme.font_family

        # Author
        if author:
            txbox3 = slide.shapes.add_textbox(
                Inches(MARGIN_LEFT), Inches(5.5),
                Inches(title_w), Inches(0.6),
            )
            tf3 = txbox3.text_frame
            tf3.word_wrap = True
            p3 = tf3.paragraphs[0]
            p3.text = author
            p3.alignment = PP_ALIGN.CENTER
            p3.font.size = resolve_font_size(self._theme, "body")
            p3.font.color.rgb = resolve_color(self._theme, "text_light")
            p3.font.name = self._theme.font_family

        sb._rendered = True  # manually built, skip auto-render
        return sb

    def section(self, num, title, subtitle=""):
        """Section divider — large 'Part N' + title."""
        sb = self.slide()
        slide = sb._slide
        title_w = SLIDE_WIDTH - 2 * MARGIN_LEFT

        # Part N
        txbox = slide.shapes.add_textbox(
            Inches(MARGIN_LEFT), Inches(2.0),
            Inches(title_w), Inches(0.8),
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Part {num}"
        p.alignment = PP_ALIGN.LEFT
        p.font.size = resolve_font_size(self._theme, "heading")
        p.font.color.rgb = resolve_color(self._theme, "primary")
        p.font.bold = True
        p.font.name = self._theme.font_family

        # Title
        txbox2 = slide.shapes.add_textbox(
            Inches(MARGIN_LEFT), Inches(2.9),
            Inches(title_w), Inches(1.2),
        )
        tf2 = txbox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title
        p2.alignment = PP_ALIGN.LEFT
        p2.font.size = Pt(36)
        p2.font.color.rgb = resolve_color(self._theme, "text")
        p2.font.bold = True
        p2.font.name = self._theme.font_family

        # Accent line
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(MARGIN_LEFT), Inches(4.3),
            Inches(1.5), Inches(0.05),
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = resolve_color(self._theme, "primary")
        line_shape.line.fill.background()

        # Subtitle
        if subtitle:
            txbox3 = slide.shapes.add_textbox(
                Inches(MARGIN_LEFT), Inches(4.6),
                Inches(title_w), Inches(0.8),
            )
            tf3 = txbox3.text_frame
            tf3.word_wrap = True
            p3 = tf3.paragraphs[0]
            p3.text = subtitle
            p3.alignment = PP_ALIGN.LEFT
            p3.font.size = resolve_font_size(self._theme, "body")
            p3.font.color.rgb = resolve_color(self._theme, "text_mid")
            p3.font.name = self._theme.font_family

        sb._rendered = True
        return sb

    def content(self, title, bullets=None, table=None, reference=None):
        """Content slide with title + bullet list or table."""
        sb = self.slide(title=title, reference=reference)

        if bullets is not None:
            sb.layout(BulletList(items=bullets, font_size="body"))
        elif table is not None:
            headers, rows = table
            sb.layout(Table(headers=headers, rows=rows))

        return sb

    def comparison(self, title, left, right, reference=None):
        """Two-column comparison slide.

        left/right: tuple of (column_title, items_list)
        """
        sb = self.slide(title=title, reference=reference)

        left_title, left_items = left
        right_title, right_items = right

        left_col = VStack(gap=0.2, children=[
            TextBlock(text=left_title, font_size="heading", bold=True,
                      color="primary", height=0.6),
            BulletList(items=left_items, font_size="body"),
        ])
        right_col = VStack(gap=0.2, children=[
            TextBlock(text=right_title, font_size="heading", bold=True,
                      color="secondary", height=0.6),
            BulletList(items=right_items, font_size="body"),
        ])

        sb.layout(HStack(gap=0.5, children=[left_col, right_col]))
        return sb

    def quote(self, text, author="", reference=None):
        """Quote slide — centered italic text with author."""
        sb = self.slide(reference=reference)
        slide = sb._slide
        title_w = SLIDE_WIDTH - 2 * MARGIN_LEFT

        # Large quote mark
        txbox_q = slide.shapes.add_textbox(
            Inches(MARGIN_LEFT + 1.0), Inches(1.5),
            Inches(1.0), Inches(1.2),
        )
        tf_q = txbox_q.text_frame
        p_q = tf_q.paragraphs[0]
        p_q.text = "\u201C"
        p_q.font.size = Pt(72)
        p_q.font.color.rgb = resolve_color(self._theme, "primary")
        p_q.font.name = self._theme.font_family

        # Quote text
        txbox = slide.shapes.add_textbox(
            Inches(MARGIN_LEFT + 1.5), Inches(2.5),
            Inches(title_w - 3.0), Inches(2.5),
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.LEFT
        p.font.size = resolve_font_size(self._theme, "heading")
        p.font.color.rgb = resolve_color(self._theme, "text")
        p.font.italic = True
        p.font.name = self._theme.font_family

        # Author
        if author:
            txbox2 = slide.shapes.add_textbox(
                Inches(MARGIN_LEFT + 1.5), Inches(5.2),
                Inches(title_w - 3.0), Inches(0.6),
            )
            tf2 = txbox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = f"-- {author}"
            p2.alignment = PP_ALIGN.RIGHT
            p2.font.size = resolve_font_size(self._theme, "body")
            p2.font.color.rgb = resolve_color(self._theme, "text_mid")
            p2.font.name = self._theme.font_family

        sb._rendered = True
        return sb

    def transition(self, text, sub_text=""):
        """Transition slide — centered bridge text."""
        sb = self.slide()
        slide = sb._slide
        title_w = SLIDE_WIDTH - 2 * MARGIN_LEFT

        txbox = slide.shapes.add_textbox(
            Inches(MARGIN_LEFT), Inches(2.5),
            Inches(title_w), Inches(1.5),
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(36)
        p.font.color.rgb = resolve_color(self._theme, "text")
        p.font.bold = True
        p.font.name = self._theme.font_family

        if sub_text:
            txbox2 = slide.shapes.add_textbox(
                Inches(MARGIN_LEFT), Inches(4.2),
                Inches(title_w), Inches(1.0),
            )
            tf2 = txbox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = sub_text
            p2.alignment = PP_ALIGN.CENTER
            p2.font.size = resolve_font_size(self._theme, "body")
            p2.font.color.rgb = resolve_color(self._theme, "text_mid")
            p2.font.name = self._theme.font_family

        sb._rendered = True
        return sb

    def end(self, title, subtitle=""):
        """Closing slide with accent line."""
        sb = self.slide()
        slide = sb._slide
        cx = SLIDE_WIDTH / 2
        title_w = SLIDE_WIDTH - 2 * MARGIN_LEFT

        # Title
        txbox = slide.shapes.add_textbox(
            Inches(MARGIN_LEFT), Inches(2.5),
            Inches(title_w), Inches(1.5),
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(40)
        p.font.color.rgb = resolve_color(self._theme, "text")
        p.font.bold = True
        p.font.name = self._theme.font_family

        # Accent line
        line_w = 2.0
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(cx - line_w / 2), Inches(4.2),
            Inches(line_w), Inches(0.05),
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = resolve_color(self._theme, "primary")
        line_shape.line.fill.background()

        if subtitle:
            txbox2 = slide.shapes.add_textbox(
                Inches(MARGIN_LEFT), Inches(4.5),
                Inches(title_w), Inches(0.8),
            )
            tf2 = txbox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.alignment = PP_ALIGN.CENTER
            p2.font.size = resolve_font_size(self._theme, "subtitle")
            p2.font.color.rgb = resolve_color(self._theme, "text_mid")
            p2.font.name = self._theme.font_family

        sb._rendered = True
        return sb

    # Bind methods to the class
    presentation_class.cover = cover
    presentation_class.section = section
    presentation_class.content = content
    presentation_class.comparison = comparison
    presentation_class.quote = quote
    presentation_class.transition = transition
    presentation_class.end = end
