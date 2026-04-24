"""PPTX backend for the IR layout pipeline."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.util import Inches, Pt

from paperops.slides.core.constants import SLIDE_HEIGHT, SLIDE_WIDTH
from paperops.slides.core.theme import Theme
from paperops.slides.ir.node import Node
from paperops.slides.layout.containers import (
    Absolute,
    AbsoluteItem,
    Flex,
    Grid,
    GridItem,
    HStack,
    LayoutNode,
    Layer,
    Padding,
)
from paperops.slides.components.shapes import Box as ShapeBox
from paperops.slides.components.table import Table as LayoutTable
from paperops.slides.components.text import TextBlock


_LITERAL_COLORS = {
    "white": RGBColor(255, 255, 255),
    "black": RGBColor(0, 0, 0),
}

_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


def render_styled_layout(
    theme: Theme,
    slide_layout_roots: list[tuple[Node, LayoutNode]],
    out_path: str | Path,
) -> Path:
    """Render laid-out nodes to PPTX."""
    out = Path(out_path)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)

    for source_slide, layout_root in slide_layout_roots:
        pptx_slide = prs.slides.add_slide(prs.slide_layouts[6])
        _render_slide_background(pptx_slide, source_slide, theme)
        _render_slide_notes(pptx_slide, source_slide)
        _render_layout_node(pptx_slide, layout_root, theme)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def _render_slide_notes(slide, source_node: Node | None) -> None:
    """Render semantic `<note>` nodes into speaker notes."""
    notes = _collect_note_texts(source_node)
    if not notes:
        return

    notes_slide = slide.notes_slide
    if notes_slide is None:
        return
    notes_slide.notes_text_frame.text = "\n\n".join(notes)


def _collect_note_texts(node: Node | None) -> list[str]:
    if node is None:
        return []
    notes: list[str] = []

    if node.type == "note":
        text = _extract_text(node, None).strip()
        if text:
            notes.append(text)

    for child in node.children or []:
        if isinstance(child, Node):
            notes.extend(_collect_note_texts(child))

    return notes


def _render_slide_background(slide, source_node: Node | None, theme: Theme) -> None:
    style = _style_snapshot(source_node)
    background = _style_value(style, "bg", None)
    if not background:
        return
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _resolve_color(theme, background)


def _render_layout_node(slide, layout_node: LayoutNode, theme: Theme) -> None:
    region = _layout_region(layout_node)
    if region is None:
        return

    source_node = getattr(layout_node, "_ir_node", None)
    if isinstance(layout_node, (Flex, HStack, Layer, Padding)):
        for child in _layout_children(layout_node):
            if child is not None:
                _render_layout_node(slide, child, theme)
        return
    if isinstance(layout_node, Absolute):
        for item in layout_node.children:
            if item is not None and item.child is not None:
                _render_layout_node(slide, item.child, theme)
        return
    if isinstance(layout_node, Grid):
        for item in layout_node.iter_items():
            if item is not None and item.child is not None:
                _render_layout_node(slide, item.child, theme)
        return
    if isinstance(layout_node, GridItem):
        if layout_node.child is not None:
            _render_layout_node(slide, layout_node.child, theme)
        return
    if isinstance(layout_node, AbsoluteItem):
        if layout_node.child is not None:
            _render_layout_node(slide, layout_node.child, theme)
        return

    _render_leaf(slide, layout_node, source_node, region, theme)


def _layout_children(node: LayoutNode) -> list[LayoutNode]:
    if isinstance(node, Padding):
        return [node.child] if node.child is not None else []
    if hasattr(node, "children"):
        children = getattr(node, "children")
        if children is None:
            return []
        return [child for child in children if child is not None]
    return []


def _layout_region(node: LayoutNode):
    region = getattr(node, "_region", None)
    if region is None:
        return None
    if region.left is None or region.top is None or region.width is None or region.height is None:
        return None
    return region


def _render_leaf(
    slide,
    layout_node: LayoutNode,
    source_node: Node | None,
    region,
    theme: Theme,
) -> None:
    text = _extract_text(source_node, layout_node)
    node_type = source_node.type if source_node is not None else None
    style = _style_snapshot(source_node)
    left = float(region.left)
    top = float(region.top)
    width = float(region.width)
    height = float(region.height)

    if node_type in {"box", "kpi"}:
        _render_box(
            slide=slide,
            region=(left, top, width, height),
            style=style,
            theme=theme,
            text=text,
        )
        return
    if node_type == "roundedbox":
        _render_rounded_box(
            slide=slide,
            region=(left, top, width, height),
            style=style,
            theme=theme,
            text=text,
        )
        return
    if node_type == "circle":
        _render_circle(
            slide=slide,
            region=(left, top, width, height),
            style=style,
            theme=theme,
            text=text,
        )
        return
    if node_type == "badge":
        _render_badge(
            slide=slide,
            region=(left, top, width, height),
            style=style,
            theme=theme,
            text=text,
        )
        return
    if node_type == "divider":
        _render_divider(
            slide=slide,
            region=(left, top, width, height),
            style=style,
            theme=theme,
            text=text,
        )
        return
    if node_type == "line":
        _render_connector(
            slide=slide,
            region=(left, top, width, height),
            style=style,
            theme=theme,
            node_type="line",
        )
        return
    if node_type == "arrow":
        _render_connector(
            slide=slide,
            region=(left, top, width, height),
            style=style,
            theme=theme,
            node_type="arrow",
        )
        return
    if node_type in {"image", "svg", "icon"}:
        if _render_image_like(slide, node_type, source_node, style, theme, (left, top, width, height), text):
            return
    if node_type == "chart":
        _render_chart(
            slide=slide,
            region=(left, top, width, height),
            style=style,
            theme=theme,
            source_node=source_node,
        )
        return

    if node_type == "table":
        if isinstance(layout_node, LayoutTable):
            _render_table(
                slide=slide,
                region=(left, top, width, height),
                style=style,
                layout_node=layout_node,
                theme=theme,
            )
        return

    if node_type == "note":
        return

    if _is_textual_node(node_type) or text:
        _render_text_box(
            slide=slide,
            region=(left, top, width, height),
            node=source_node,
            style=style,
            theme=theme,
            text=text,
        )


def _is_textual_node(node_type: str | None) -> bool:
    if node_type is None:
        return False
    return node_type in {"text", "title", "subtitle", "heading"}


def _node_type_defaults(node_type: str | None) -> str:
    match node_type:
        case "title":
            return "title"
        case "subtitle":
            return "subtitle"
        case "heading":
            return "heading"
        case _:
            return "body"


def _text_color(style: dict[str, Any], theme: Theme) -> RGBColor:
    return _resolve_color(theme, _style_value(style, "color", "text"))


def _render_box(
    *,
    slide,
    region: tuple[float, float, float, float],
    style: dict[str, Any],
    theme: Theme,
    text: str,
) -> None:
    left, top, width, height = region
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )

    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = _resolve_color(
        theme, _style_value(style, "bg", _theme_default_color(theme, "bg_alt"))
    )

    if _style_value(style, "border", "none") not in {"none", None, "inherit"}:
        shape.line.color.rgb = _resolve_color(
            theme,
            _style_value(style, "border", _theme_default_color(theme, "border")),
        )
        line_width = _to_float(_style_value(style, "line-width", 1.0))
        shape.line.width = Pt(line_width if line_width is not None else 1.0)
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        _apply_text_frame_style(tf, style, theme=theme)
        paragraph = tf.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = _ALIGN_MAP.get(
            _style_value(style, "text-align", "center"), PP_ALIGN.LEFT
        )


def _theme_default_color(theme: Theme, key: str) -> str:
    return str(theme.colors.get(key, "#FFFFFF"))


def _render_rounded_box(
    *,
    slide,
    region: tuple[float, float, float, float],
    style: dict[str, Any],
    theme: Theme,
    text: str,
) -> None:
    left, top, width, height = region
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )

    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = _resolve_color(
        theme, _style_value(style, "bg", _theme_default_color(theme, "bg_alt"))
    )

    shape.line.fill.background()
    if _style_value(style, "border", "none") not in {"none", None, "inherit"}:
        shape.line.color.rgb = _resolve_color(
            theme,
            _style_value(style, "border", _theme_default_color(theme, "border")),
        )
        line_width = _to_float(_style_value(style, "line-width", 1.0))
        shape.line.width = Pt(line_width if line_width is not None else 1.0)

    if text:
        tf = shape.text_frame
        _apply_text_frame_style(tf, style, theme=theme)
        paragraph = tf.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = _ALIGN_MAP.get(
            _style_value(style, "text-align", "center"), PP_ALIGN.LEFT
        )


def _render_circle(
    *,
    slide,
    region: tuple[float, float, float, float],
    style: dict[str, Any],
    theme: Theme,
    text: str,
) -> None:
    left, top, width, height = region
    diameter = max(width, height)
    cx = left + (width - diameter) / 2.0
    cy = top + (height - diameter) / 2.0
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx),
        Inches(cy),
        Inches(diameter),
        Inches(diameter),
    )

    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = _resolve_color(
        theme, _style_value(style, "bg", _theme_default_color(theme, "primary"))
    )
    shape.line.fill.background()

    if text:
        tf = shape.text_frame
        _apply_text_frame_style(tf, style, theme=theme)
        paragraph = tf.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = PP_ALIGN.CENTER


def _render_badge(
    *,
    slide,
    region: tuple[float, float, float, float],
    style: dict[str, Any],
    theme: Theme,
    text: str,
) -> None:
    left, top, width, height = region
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = _resolve_color(
        theme, _style_value(style, "bg", _theme_default_color(theme, "accent"))
    )
    shape.line.fill.background()

    if text:
        tf = shape.text_frame
        _apply_text_frame_style(tf, style, theme=theme)
        paragraph = tf.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = PP_ALIGN.CENTER


def _render_divider(
    *,
    slide,
    region: tuple[float, float, float, float],
    style: dict[str, Any],
    theme: Theme,
    text: str,
)-> None:
    style_orientation = _style_value(style, "orientation", _style_value(style, "direction", "horizontal"))
    is_vertical = str(style_orientation).lower() in {"vertical", "v"}

    left, top, width, height = region
    if is_vertical:
        x1 = Inches(left + width / 2.0)
        y1 = Inches(top)
        x2 = Inches(left + width / 2.0)
        y2 = Inches(top + height)
    else:
        x1 = Inches(left)
        y1 = Inches(top + height / 2.0)
        x2 = Inches(left + width)
        y2 = Inches(top + height / 2.0)

    shape = slide.shapes.add_connector(1, x1, y1, x2, y2)
    shape.line.color.rgb = _resolve_color(
        theme, _style_value(style, "color", _theme_default_color(theme, "border"))
    )
    line_width = _to_float(_style_value(style, "line-width", 1.0))
    shape.line.width = Pt(line_width if line_width is not None else 1.0)


def _render_connector(
    *,
    slide,
    region: tuple[float, float, float, float],
    style: dict[str, Any],
    theme: Theme,
    node_type: str,
) -> None:
    left, top, width, height = region
    is_vertical = str(_style_value(style, "orientation", "horizontal")).lower() in {"vertical", "v"}

    x1 = left + 0.0
    y1 = top + height / 2.0
    x2 = left + width
    y2 = top + height / 2.0
    if is_vertical:
        x1 = left + width / 2.0
        y1 = top
        x2 = left + width / 2.0
        y2 = top + height

    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = _resolve_color(
        theme, _style_value(style, "color", _theme_default_color(theme, "border"))
    )
    line_width = _to_float(_style_value(style, "line-width", 1.0))
    connector.line.width = Pt(line_width if line_width is not None else 1.0)

    if node_type == "arrow" and width > 0.0 and height > 0.0:
        line_color = _resolve_color(
            theme, _style_value(style, "color", _theme_default_color(theme, "border"))
        )
        connector.line.color.rgb = line_color
        from pptx.oxml.ns import qn

        line_element = connector.line._ln
        arrow_tip = line_element.makeelement(qn("a:headEnd"), {})
        arrow_tip.set("type", "triangle")
        arrow_tip.set("w", "med")
        arrow_tip.set("len", "med")
        line_element.append(arrow_tip)


def _render_image_like(
    slide,
    node_type: str,
    source_node: Node | None,
    style: dict[str, Any],
    theme: Theme,
    region: tuple[float, float, float, float],
    extracted_text: str,
) -> bool:
    props = source_node.props if source_node is not None else {}
    left, top, width, height = region
    if node_type == "icon":
        text = props.get("name")
        if text:
            _render_text_box(
                slide=slide,
                region=(left, top, width, height),
                node=source_node,
                style=style,
                theme=theme,
                text=f"[icon] {text}",
            )
            return True
        if extracted_text:
            _render_text_box(
                slide=slide,
                region=(left, top, width, height),
                node=source_node,
                style=style,
                theme=theme,
                text=extracted_text,
            )
            return True
        return False

    src = ""
    if isinstance(props, dict):
        src = str(props.get("src", "") or "")

    if node_type == "svg" and not src:
        src = str(props.get("body", "") or "")
        if src.startswith("<") and src.endswith(">"):
            placeholder = extracted_text or "svg"
            _render_text_box(
                slide=slide,
                region=(left, top, width, height),
                node=source_node,
                style=style,
                theme=theme,
                text=placeholder,
            )
            return True
        return False

    if src and src != "missing.png" and Path(src).is_file():
        path = Path(src)
        try:
            slide.shapes.add_picture(
                str(path),
                Inches(left),
                Inches(top),
                Inches(width),
                Inches(height),
            )
            return True
        except Exception:
            pass

    label = extracted_text or node_type
    _render_text_box(
        slide=slide,
        region=(left, top, width, height),
        node=source_node,
        style=style,
        theme=theme,
        text=label,
    )
    return True


def _render_table(
    slide,
    region: tuple[float, float, float, float],
    style: dict[str, Any],
    layout_node: LayoutTable,
    theme: Theme,
) -> None:
    left, top, width, height = region

    rows = list(getattr(layout_node, "rows", []))
    headers = list(getattr(layout_node, "headers", []))

    if not rows and not headers:
        rows = [[""]]

    row_count = len(rows) + (1 if headers else 0)
    row_count = max(row_count, 1)

    column_count = max(len(headers), 1)
    for row in rows:
        column_count = max(column_count, len(row))

    table_shape = slide.shapes.add_table(
        row_count,
        column_count,
        Inches(left),
        Inches(top),
        Inches(max(width, 0.1)),
        Inches(max(height, 0.1)),
    )
    table = table_shape.table

    if headers:
        header_row = table.rows[0]
        for column_index, value in enumerate(headers):
            if column_index >= len(header_row.cells):
                break
            header_cell = header_row.cells[column_index]
            header_cell.text = _as_text_safe(value)
            paragraph = header_cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.name = theme.font_family
            paragraph.font.color.rgb = _text_color(style, theme)
            paragraph.alignment = PP_ALIGN.CENTER
        start_row = 1
        if len(table.rows) > 0:
            table.rows[0].height = Pt(_resolve_font_size(theme, "body", "body"))
    else:
        start_row = 0

    for row_offset, row_data in enumerate(rows):
        if start_row + row_offset >= len(table.rows):
            break
        target_row = table.rows[start_row + row_offset]
        for column_index, value in enumerate(row_data):
            if column_index >= len(target_row.cells):
                break
            target_cell = target_row.cells[column_index]
            target_cell.text = _as_text_safe(value)

    cell_bg = _style_value(style, "cell-bg")
    if cell_bg is not None:
        try:
            bg = _resolve_color(theme, cell_bg)
        except Exception:
            bg = _resolve_color(theme, "bg")
        for row in table.rows:
            for cell in row.cells:
                if getattr(cell, "fill", None) is not None:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = bg


def _as_text_safe(value: Any) -> str:
    if value is None:
        return ""
    return str(value)


def _render_text_box(
    *,
    slide,
    region: tuple,
    node: Node | None,
    style: dict[str, Any],
    theme: Theme,
    text: str,
) -> None:
    left, top, width, height = region
    tb = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    tf = tb.text_frame
    _apply_text_frame_style(tf, style, theme=theme)
    paragraph = tf.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = _ALIGN_MAP.get(_style_value(style, "text-align", "left"), PP_ALIGN.LEFT)

    node_type = node.type if node is not None else None
    font_name = theme.font_family
    font_token = _style_value(style, "font", _node_type_defaults(node_type))
    paragraph.font.name = font_name
    paragraph.font.size = Pt(_resolve_font_size(theme, font_token, "body"))
    paragraph.font.color.rgb = _text_color(style, theme)
    paragraph.font.bold = _as_bool(_style_value(style, "font-weight", "normal"))
    paragraph.font.italic = str(_style_value(style, "font-style", "")).lower() == "italic"
    if (line_height := _style_value(style, "line-height")) is not None:
        numeric = _to_float(line_height)
        if numeric is not None:
            paragraph.line_spacing = numeric


def _chart_type_to_xl(chart_type: str) -> XL_CHART_TYPE:
    normalized = str(chart_type or "").strip().lower()
    mapping = {
        "line": XL_CHART_TYPE.LINE,
        "line_markers": XL_CHART_TYPE.LINE_MARKERS,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
        "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
        "pie": XL_CHART_TYPE.PIE,
        "area": XL_CHART_TYPE.AREA,
        "stacked_area": XL_CHART_TYPE.AREA_STACKED,
    }
    return mapping.get(normalized, XL_CHART_TYPE.LINE)


def _coerce_chart_values(raw_values: object) -> list[float]:
    if isinstance(raw_values, (int, float, str)):
        return [_safe_float(raw_values)]
    if not isinstance(raw_values, (list, tuple)):
        return []
    values: list[float] = []
    for raw_value in raw_values:
        value = _safe_float(raw_value)
        if value is not None:
            values.append(value)
    return values


def _safe_float(value: object) -> float | None:
    if value is None:
        return None
    if isinstance(value, (int, float)):
        return float(value)
    try:
        numeric = float(str(value))
        return numeric
    except (TypeError, ValueError):
        return None


def _coerce_chart_series(payload: object) -> list[tuple[str, list[float]]]:
    if payload is None:
        return []

    if isinstance(payload, dict):
        series: list[tuple[str, list[float]]] = []
        for name, values in payload.items():
            values_list = _coerce_chart_values(values)
            if values_list:
                series.append((str(name), values_list))
        return series

    if not isinstance(payload, (list, tuple)):
        return []

    series: list[tuple[str, list[float]]] = []
    for item in payload:
        if isinstance(item, dict):
            name = item.get("name") or item.get("label") or item.get("title")
            if "data" in item and name is None:
                name = "Series"
            values = item.get("values") if isinstance(item, dict) else None
            if name is None:
                name = f"Series {len(series) + 1}"
            values_list = _coerce_chart_values(values if values is not None else item.get("data", []))
            if values_list:
                series.append((str(name), values_list))
            continue
        if isinstance(item, (list, tuple)):
            values = [value for value in item if _safe_float(value) is not None]
            if values:
                series.append((f"Series {len(series) + 1}", _coerce_chart_values(values)))
    return series


def _resolve_chart_payload(
    source_node: Node | None,
    props: dict[str, object],
) -> tuple[str, list[str], list[tuple[str, list[float]]]]:
    chart_type = ""
    if isinstance(props.get("chart_type"), str):
        chart_type = str(props["chart_type"])

    labels: list[str] | None = None
    raw_labels = props.get("labels")
    if isinstance(raw_labels, (list, tuple)):
        labels = [str(label) for label in raw_labels]

    series: list[tuple[str, list[float]]] = []
    if isinstance(props.get("series"), (list, tuple, dict)):
        series = _coerce_chart_series(props["series"])

    data = props.get("data")
    if isinstance(data, dict):
        if not series and isinstance(data.get("series"), (list, tuple, dict)):
            series = _coerce_chart_series(data.get("series"))

        if not labels and "labels" in data and isinstance(data["labels"], (list, tuple)):
            labels = [str(label) for label in data["labels"]]

        if isinstance(data.get("values"), (list, tuple)):
            values = _coerce_chart_values(data["values"])
            if values and not series:
                series = [("Series", values)]

        if not series:
            flat_data: list[float] = []
            flat_labels: list[str] = []
            for key, value in data.items():
                maybe_value = _safe_float(value)
                if maybe_value is not None:
                    flat_labels.append(str(key))
                    flat_data.append(maybe_value)
            if flat_data:
                if labels is None:
                    labels = flat_labels
                series = [("Series", flat_data)]

    if not labels and source_node is not None and source_node.text:
        labels = [str(source_node.text)]

    if not series:
        series = [("Series", [0.0])]
        if labels is None:
            labels = ["Point 1"]
    elif labels is None and series:
        labels = [f"Point {idx + 1}" for idx in range(max(len(values) for _, values in series))]

    normalized_labels: list[str] = []
    if labels:
        normalized_labels = [str(item) for item in labels]

    max_len = max(len(values) for _, values in series) if series else 0
    if normalized_labels and len(normalized_labels) < max_len:
        normalized_labels.extend(f"Point {idx + 1}" for idx in range(len(normalized_labels) + 1, max_len + 1))
    elif not normalized_labels and max_len:
        normalized_labels = [f"Point {idx + 1}" for idx in range(1, max_len + 1)]

    return chart_type, normalized_labels, series


def _render_chart(
    *,
    slide,
    region: tuple[float, float, float, float],
    style: dict[str, Any],
    theme: Theme,
    source_node: Node | None,
) -> None:
    left, top, width, height = region
    props = source_node.props if source_node is not None and source_node.props is not None else {}
    if not isinstance(props, dict):
        props = {}

    chart_type, labels, series = _resolve_chart_payload(source_node, props)
    title = _style_value(style, "title", None)
    if title is None and source_node is not None and source_node.text:
        title = source_node.text
    if not props.get("title"):
        props["title"] = title

    chart_data = ChartData()
    chart_data.categories = list(labels)
    for name, values in series:
        chart_data.add_series(str(name), values)

    try:
        chart_shape = slide.shapes.add_chart(
            _chart_type_to_xl(chart_type),
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
            chart_data,
        )
    except TypeError:
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
            chart_data,
        )

    chart = chart_shape.chart
    if props.get("title") or props.get("caption"):
        chart.has_title = True
        caption = props.get("title") or props.get("caption")
        chart.chart_title.text_frame.text = str(caption)

    if len(series) <= 1:
        chart.has_legend = False

    if not chart_data.categories:
        return

    for plot in chart.plots:
        for series_idx in range(len(plot.series)):
            plot.series[series_idx].has_data_labels = False


def _apply_text_frame_style(text_frame, style: dict[str, Any], *, theme: Theme) -> None:
    font_name = theme.font_family
    text_frame.margin_left = Inches(_padding_or_margin(style, "left"))
    text_frame.margin_right = Inches(_padding_or_margin(style, "right"))
    text_frame.margin_top = Inches(_padding_or_margin(style, "top"))
    text_frame.margin_bottom = Inches(_padding_or_margin(style, "bottom"))

    text_frame.word_wrap = True
    if not text_frame.text:
        first = text_frame.paragraphs[0]
        first.font.name = font_name
        first.font.size = Pt(_resolve_font_size(theme, _style_value(style, "font", "body"), "body"))
        first.font.color.rgb = _text_color(style, theme)


def _padding_or_margin(style: dict[str, Any], side: str) -> float:
    if style is None:
        return 0.0
    keys = {
        "left": ("padding-left", "padding-x", "padding"),
        "right": ("padding-right", "padding-x", "padding"),
        "top": ("padding-top", "padding-y", "padding"),
        "bottom": ("padding-bottom", "padding-y", "padding"),
    }[side]
    for key in keys:
        value = style.get(key)
        if value is None:
            continue
        parsed = _to_float(value)
        if parsed is not None:
            return parsed
    return 0.0


def _resolve_font_size(theme: Theme, value: Any, fallback: str) -> float:
    if isinstance(value, (int, float)):
        return float(value)
    token = value if value is not None else fallback
    try:
        return float(theme.resolve_font_size(token))
    except Exception:
        return float(theme.resolve_font_size(fallback))


def _style_value(style: dict[str, Any], key: str, default: Any = None) -> Any:
    if not style:
        return default
    value = style.get(key, default)
    if value in {"auto", "inherit", None}:
        return default
    return value


def _resolve_color(theme: Theme, value: Any) -> RGBColor:
    if isinstance(value, RGBColor):
        return value
    if not isinstance(value, str):
        raise TypeError(f"Unsupported color value: {value!r}")
    if value in _LITERAL_COLORS:
        return _LITERAL_COLORS[value]
    resolved = theme.resolve_color(value)
    if not isinstance(resolved, str):
        raise TypeError(f"Unsupported color token value: {resolved!r}")
    if not resolved.startswith("#") or len(resolved) != 7:
        raise ValueError(f"Unsupported color format: {resolved!r}")
    return RGBColor(int(resolved[1:3], 16), int(resolved[3:5], 16), int(resolved[5:7], 16))


def _as_bool(value: Any) -> bool:
    if value is None:
        return False
    if isinstance(value, bool):
        return value
    if isinstance(value, (int, float)):
        return value >= 1
    if isinstance(value, str):
        return value.lower() in {"1", "true", "yes", "bold", "bolder"}
    return False


def _to_float(value: Any) -> float | None:
    if value is None or value in {"auto", "inherit"}:
        return None
    if isinstance(value, (int, float)):
        return float(value)
    try:
        return float(str(value))
    except (TypeError, ValueError):
        return None


def _extract_text(source_node: Node | None, layout_node: LayoutNode | None) -> str:
    if source_node is not None:
        if source_node.text is not None:
            return str(source_node.text)
        if source_node.props and isinstance(source_node.props.get("text"), str):
            return str(source_node.props["text"])
        if source_node.children:
            chunks: list[str] = []
            for child in source_node.children:
                if isinstance(child, str):
                    chunks.append(child)
                elif isinstance(child, Node):
                    chunks.append(_extract_text(child, None))
            return "".join(chunks)

    if isinstance(layout_node, TextBlock):
        return str(layout_node.text)
    if isinstance(layout_node, ShapeBox):
        return str(layout_node.text)
    return ""


def _style_snapshot(node: Node | None) -> dict[str, Any]:
    if node is None:
        return {}
    style = node.style or {}
    if hasattr(node, "computed_style"):
        snapshot = node.computed_style.snapshot()
        if snapshot:
            style = {**snapshot, **style}
    return style
