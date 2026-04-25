"""Pipeline entrypoint for IR-based deck rendering and legacy Presentation API."""

from __future__ import annotations

from pathlib import Path
from typing import Any, Mapping

import os
import sys
import tempfile

from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from paperops.slides import components  # noqa: F401
from paperops.slides.codegen import render_styled_layout
from paperops.slides.core.constants import CONTENT_REGION, SLIDE_HEIGHT, SLIDE_WIDTH
from paperops.slides.core.theme import Theme, themes
from paperops.slides.dsl.json_loader import Document
from paperops.slides.dsl.json_loader import load_json_document
from paperops.slides.dsl.markdown_parser import load_markdown_document
from paperops.slides.dsl.mdx_parser import load_mdx_document
from paperops.slides.layout.engine import build_layout_tree, compute_layout
from paperops.slides.components.registry import expand_nodes
from paperops.slides.style import get_sheet, resolve_computed_styles
from paperops.slides.style.stylesheet import StyleSheet
from paperops.slides.preview import review_deck_artifacts

from paperops.slides.ir.node import Node


def _coerce_theme(value: Theme | str | None) -> Theme:
    if value is None:
        return themes.minimal
    if isinstance(value, Theme):
        return value
    if isinstance(value, str):
        try:
            candidate = getattr(themes, value)
        except AttributeError as exc:
            raise ValueError(f"Unknown theme: {value!r}") from exc
        if not isinstance(candidate, Theme):
            raise TypeError(f"Theme {value!r} did not resolve to a Theme instance")
        return candidate
    raise TypeError(f"theme must be a Theme or string, got {type(value)!r}")


def parse_stage(
    source: str | Path | Mapping[str, Any] | Document,
    *,
    strict: bool = False,
) -> Document:
    """Parse JSON or IR document into a canonical :class:`Document`."""
    if isinstance(source, Document):
        return source
    if isinstance(source, Mapping):
        return load_json_document(dict(source), strict=strict)

    path = Path(source)
    if path.exists() and path.is_file():
        name = path.name.lower()
        if name.endswith(".deck.md") or name.endswith(".md"):
            return load_markdown_document(path)
        if name.endswith(".slide.mdx") or name.endswith(".deck.mdx") or name.endswith(".mdx"):
            return load_mdx_document(path)

    return load_json_document(source, strict=strict)


def expand_stage(document: Document, *, strict: bool = False) -> list[Node]:
    """Expand semantic components and return expanded slide roots."""
    return expand_nodes(document.slides, strict=strict)


def style_stage(
    document: Document,
    *,
    theme: Theme | str | None = None,
    sheet: Mapping[str, Mapping[str, Any]] | str | None = None,
    strict: bool = False,
) -> tuple[Node, dict[int, Any]]:
    """Resolve computed styles and attach them to every styled node."""
    if sheet is None and document.sheet is not None:
        sheet = document.sheet
    if isinstance(sheet, str):
        sheet = get_sheet(sheet)

    style_theme = _coerce_theme(theme or document.theme)
    expanded_slides = expand_stage(document, strict=strict)
    root = Node(type="deck", children=expanded_slides)
    resolved = resolve_computed_styles(
        root,
        theme=style_theme,
        style_sheet=sheet,
        deck_style=StyleSheet(document.styles or {}),
        strict=strict,
    )
    return root, resolved.computed


def layout_stage(
    styled_root: Node,
    *,
    theme: Theme | str | None = None,
    region=CONTENT_REGION,
    slide: int | None = None,
) -> tuple[list[tuple[Node, Any]], list[dict[str, Any]]]:
    """Compute regions for each slide node."""
    layout_theme = _coerce_theme(theme)
    slide_layouts: list[tuple[Node, Any]] = []
    issues: list[dict[str, Any]] = []

    for index, slide_node in enumerate(styled_root.children or []):
        if not isinstance(slide_node, Node):
            continue
        layout_root = build_layout_tree(slide_node, layout_theme, region=region)
        slide_layouts.append((slide_node, layout_root))
        issues.extend(
            compute_layout(layout_root, region, layout_theme, slide=index + 1, root_path=f"slide[{index}]")
        )
    return slide_layouts, issues


def autofit_stage(
    slide_layouts: list[tuple[Node, Any]],
    *,
    theme: Theme | str | None = None,
) -> list[tuple[Node, Any]]:
    """Placeholder autofit stage for Phase 2 (currently identity)."""
    return slide_layouts


def codegen_stage(
    slide_layouts: list[tuple[Node, Any]],
    *,
    theme: Theme | str | None = None,
    out: str | Path,
) -> Path:
    resolved_theme = _coerce_theme(theme)
    return render_styled_layout(resolved_theme, slide_layouts, out_path=out)


def render_json(
    source: str | Path | Mapping[str, Any] | Document,
    *,
    out: str | Path,
    strict: bool = False,
) -> Path:
    """Run parse → style → layout → codegen and write a PPTX file."""
    document = parse_stage(source, strict=strict)
    styled_root, _ = style_stage(document, theme=document.theme, strict=strict)
    layout_roots = autofit_stage(layout_stage(styled_root, theme=document.theme)[0], theme=document.theme)
    return codegen_stage(layout_roots, theme=document.theme, out=out)


def _normalize_issue(issue: dict, slide_number: int) -> dict[str, Any]:
    normalized = dict(issue)
    normalized.setdefault("slide", slide_number)
    normalized.setdefault("source", "layout")
    normalized.setdefault("code", normalized.get("type", "layout_issue"))
    normalized.setdefault("message", normalized.get("detail", normalized.get("code", "layout issue")))
    normalized.setdefault("detail", normalized["message"])
    normalized.setdefault("severity", "warning")
    return normalized


# ---------------------------------------------------------------------------
# Legacy API compatibility
# ---------------------------------------------------------------------------


class Presentation:
    """Main entry point for the original imperative component builder API."""

    def __init__(self, theme: Theme | None = None):
        self._theme = theme if theme is not None else themes.professional
        self._pptx = PptxPresentation()
        self._pptx.slide_width = Inches(SLIDE_WIDTH)
        self._pptx.slide_height = Inches(SLIDE_HEIGHT)
        self._builders: list[object] = []

    def slide(self, title=None, reference=None, background=None):
        from paperops.slides.slides.base import SlideBuilder

        layout = self._pptx.slide_layouts[6]
        pptx_slide = self._pptx.slides.add_slide(layout)
        sb = SlideBuilder(
            pptx_slide,
            self._theme,
            title=title,
            reference=reference,
            slide_number=len(self._builders) + 1,
        )
        if background is not None:
            sb.background(color=background)
        self._builders.append(sb)
        return sb

    def save(self, path: str):
        all_issues = []
        for slide_index, sb in enumerate(self._builders, 1):
            issues = sb._render(slide_number=slide_index)
            if issues:
                all_issues.extend(issues)

        if all_issues:
            print(f"[SlideCraft] {len(all_issues)} layout issue(s):", file=sys.stderr)
            for issue in all_issues:
                print(f"  [{issue.get('code', issue.get('type'))}] {issue.get('detail', issue.get('message'))}", file=sys.stderr)

        self._pptx.save(path)

    def review(self) -> dict:
        issues: list[dict] = []
        slides: list[dict] = []
        for slide_index, sb in enumerate(self._builders, 1):
            slide_issues = sb._render(slide_number=slide_index) or []
            normalized = [_normalize_issue(issue, slide_index) for issue in slide_issues]
            issues.extend(normalized)
            slides.append({
                "slide_number": slide_index,
                "title": getattr(sb, "_title", None),
                "issue_count": len(normalized),
                "issues": normalized,
            })

        issue_counts = {
            "total": len(issues),
            "error": sum(issue.get("severity") == "error" for issue in issues),
            "warning": sum(issue.get("severity") == "warning" for issue in issues),
            "info": sum(issue.get("severity") == "info" for issue in issues),
        }
        return {
            "schema_version": "2026-04-09",
            "artifact": "layout_review",
            "total_slides": len(self._builders),
            "total_issues": len(issues),
            "issue_counts": issue_counts,
            "issues": issues,
            "slides": slides,
        }

    def review_deck(
        self,
        output_path: str,
        render_preview: bool = True,
        output_dir: str | None = None,
    ) -> dict:
        layout_review = self.review()
        self.save(output_path)
        preview_paths: list[str] = []
        if render_preview:
            if output_dir is not None:
                os.makedirs(output_dir, exist_ok=True)
                for filename in os.listdir(output_dir):
                    if filename.startswith("slide_") and filename.endswith(".png"):
                        os.unlink(os.path.join(output_dir, filename))
            preview_paths = self.preview(output_dir=output_dir)
        slide_titles = [getattr(sb, "_title", None) for sb in self._builders]
        return review_deck_artifacts(
            output_path,
            layout_issues=layout_review["issues"],
            preview_paths=preview_paths,
            slide_titles=slide_titles,
            layout_summary=layout_review,
        )

    def preview(self, slides=None, output_dir=None) -> list[str]:
        from paperops.slides.preview import render_slide_preview_powerpoint

        if output_dir is None:
            output_dir = tempfile.mkdtemp(prefix="slidecraft_preview_")
        os.makedirs(output_dir, exist_ok=True)

        for slide_index, sb in enumerate(self._builders, 1):
            sb._render(slide_number=slide_index)

        pptx_path = os.path.join(output_dir, "_temp_preview.pptx")
        self.save(pptx_path)
        paths = render_slide_preview_powerpoint(pptx_path, output_dir)
        os.unlink(pptx_path)

        if slides is not None:
            paths = [p for p in paths if any(f"slide_{s + 1:03d}.png" in p for s in slides)]
        return sorted(paths)


__all__ = [
    "Presentation",
    "autofit_stage",
    "codegen_stage",
    "layout_stage",
    "parse_stage",
    "render_json",
    "style_stage",
]
