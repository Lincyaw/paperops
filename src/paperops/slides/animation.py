"""Animation system -- generates OOXML timing XML for click-to-advance animations."""

from lxml import etree

NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}

_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'


def _pn(tag: str) -> str:
    """Return a fully qualified tag name in the PresentationML namespace."""
    return f'{{{_P}}}{tag}'


def _inject_initial_hide(parent_child_tnlst, shape_ids: list[int], id_start: int) -> int:
    """Add an auto-playing animation at slide start that hides shapes.

    Creates a <p:par> that triggers at delay=0 (auto) and sets
    style.visibility=hidden for each shape in shape_ids.

    Returns the next available ID counter after all IDs allocated here.
    """
    if not shape_ids:
        return id_start

    hide_par = etree.SubElement(parent_child_tnlst, _pn('par'))
    hide_cTn = etree.SubElement(hide_par, _pn('cTn'))
    hide_cTn.set('id', str(id_start))
    hide_cTn.set('fill', 'hold')

    # Trigger immediately when slide loads
    stCond = etree.SubElement(hide_cTn, _pn('stCondLst'))
    cond = etree.SubElement(stCond, _pn('cond'))
    cond.set('delay', '0')

    hide_child = etree.SubElement(hide_cTn, _pn('childTnLst'))

    id_counter = id_start + 1

    for spid_int in shape_ids:
        spid = str(spid_int)

        shape_par = etree.SubElement(hide_child, _pn('par'))
        shape_cTn = etree.SubElement(shape_par, _pn('cTn'))
        shape_cTn.set('id', str(id_counter)); id_counter += 1
        shape_cTn.set('fill', 'hold')

        shape_stCond = etree.SubElement(shape_cTn, _pn('stCondLst'))
        shape_condEl = etree.SubElement(shape_stCond, _pn('cond'))
        shape_condEl.set('delay', '0')

        shape_child = etree.SubElement(shape_cTn, _pn('childTnLst'))

        # Set style.visibility = hidden
        setEl = etree.SubElement(shape_child, _pn('set'))

        set_cBhvr = etree.SubElement(setEl, _pn('cBhvr'))
        set_cTn = etree.SubElement(set_cBhvr, _pn('cTn'))
        set_cTn.set('id', str(id_counter)); id_counter += 1
        set_cTn.set('dur', '1')
        set_cTn.set('fill', 'hold')

        set_stCond = etree.SubElement(set_cTn, _pn('stCondLst'))
        set_condEl = etree.SubElement(set_stCond, _pn('cond'))
        set_condEl.set('delay', '0')

        set_tgtEl = etree.SubElement(set_cBhvr, _pn('tgtEl'))
        set_spTgt = etree.SubElement(set_tgtEl, _pn('spTgt'))
        set_spTgt.set('spid', spid)

        set_attrLst = etree.SubElement(set_cBhvr, _pn('attrNameLst'))
        set_attr = etree.SubElement(set_attrLst, _pn('attrName'))
        set_attr.text = 'style.visibility'

        set_to = etree.SubElement(setEl, _pn('to'))
        set_strVal = etree.SubElement(set_to, _pn('strVal'))
        set_strVal.set('val', 'hidden')

    return id_counter


def inject_appear_animations(slide_element, click_groups: list[list[int]], *,
                              initially_hidden: list[int] | None = None):
    """Inject click-to-advance fade animations into a slide's XML.

    Args:
        slide_element: the lxml element for the slide (slide._element in python-pptx)
        click_groups: list of lists of shape IDs (int).
            Each inner list = shapes that appear together on one click.
            Groups appear in order: click 1 reveals group[0], click 2 reveals group[1], etc.
        initially_hidden: shape IDs that should start invisible (typically groups 2+).
    """
    if not click_groups:
        return

    # Root timing structure
    timing = etree.SubElement(slide_element, _pn('timing'))
    tnLst = etree.SubElement(timing, _pn('tnLst'))
    root_par = etree.SubElement(tnLst, _pn('par'))

    root_cTn = etree.SubElement(root_par, _pn('cTn'))
    root_cTn.set('id', '1')
    root_cTn.set('dur', 'indefinite')
    root_cTn.set('restart', 'never')
    root_cTn.set('nodeType', 'tmRoot')

    root_child = etree.SubElement(root_cTn, _pn('childTnLst'))

    # --- Initial hide: auto-playing animation that hides shapes at slide start ---
    if initially_hidden:
        id_counter = _inject_initial_hide(root_child, initially_hidden, id_start=3)
    else:
        id_counter = 3

    seq = etree.SubElement(root_child, _pn('seq'))
    seq.set('concurrent', '1')
    seq.set('nextAc', 'seek')

    seq_cTn = etree.SubElement(seq, _pn('cTn'))
    seq_cTn.set('id', str(id_counter)); id_counter += 1
    seq_cTn.set('dur', 'indefinite')
    seq_cTn.set('nodeType', 'mainSeq')

    seq_child = etree.SubElement(seq_cTn, _pn('childTnLst'))

    for group_idx, group in enumerate(click_groups):
        # Click group container
        click_par = etree.SubElement(seq_child, _pn('par'))
        click_cTn = etree.SubElement(click_par, _pn('cTn'))
        click_cTn.set('id', str(id_counter)); id_counter += 1
        click_cTn.set('fill', 'hold')

        stCond = etree.SubElement(click_cTn, _pn('stCondLst'))
        cond = etree.SubElement(stCond, _pn('cond'))
        cond.set('delay', 'indefinite')  # wait for mouse click

        click_child = etree.SubElement(click_cTn, _pn('childTnLst'))

        # Inner par for simultaneous playback
        inner_par = etree.SubElement(click_child, _pn('par'))
        inner_cTn = etree.SubElement(inner_par, _pn('cTn'))
        inner_cTn.set('id', str(id_counter)); id_counter += 1
        inner_cTn.set('fill', 'hold')

        inner_stCond = etree.SubElement(inner_cTn, _pn('stCondLst'))
        inner_condEl = etree.SubElement(inner_stCond, _pn('cond'))
        inner_condEl.set('delay', '0')

        inner_child = etree.SubElement(inner_cTn, _pn('childTnLst'))

        for shape_idx, spid_int in enumerate(group):
            spid = str(spid_int)

            # Animation node for this shape
            anim_par = etree.SubElement(inner_child, _pn('par'))
            anim_cTn = etree.SubElement(anim_par, _pn('cTn'))
            anim_cTn.set('id', str(id_counter)); id_counter += 1
            anim_cTn.set('presetID', '10')       # Fade
            anim_cTn.set('presetClass', 'entr')
            anim_cTn.set('presetSubtype', '0')
            anim_cTn.set('fill', 'hold')
            anim_cTn.set('grpId', '0')
            if shape_idx == 0:
                anim_cTn.set('nodeType', 'clickEffect')
            else:
                anim_cTn.set('nodeType', 'withEffect')

            anim_stCond = etree.SubElement(anim_cTn, _pn('stCondLst'))
            anim_condEl = etree.SubElement(anim_stCond, _pn('cond'))
            anim_condEl.set('delay', '0')

            anim_child = etree.SubElement(anim_cTn, _pn('childTnLst'))

            # --- Fade effect ---
            animEffect = etree.SubElement(anim_child, _pn('animEffect'))
            animEffect.set('transition', 'in')
            animEffect.set('filter', 'fade')

            eff_cBhvr = etree.SubElement(animEffect, _pn('cBhvr'))
            eff_cTn = etree.SubElement(eff_cBhvr, _pn('cTn'))
            eff_cTn.set('id', str(id_counter)); id_counter += 1
            eff_cTn.set('dur', '500')

            eff_tgtEl = etree.SubElement(eff_cBhvr, _pn('tgtEl'))
            eff_spTgt = etree.SubElement(eff_tgtEl, _pn('spTgt'))
            eff_spTgt.set('spid', spid)

            # --- Visibility set ---
            setEl = etree.SubElement(anim_child, _pn('set'))

            set_cBhvr = etree.SubElement(setEl, _pn('cBhvr'))
            set_cTn = etree.SubElement(set_cBhvr, _pn('cTn'))
            set_cTn.set('id', str(id_counter)); id_counter += 1
            set_cTn.set('dur', '1')
            set_cTn.set('fill', 'hold')

            set_stCond = etree.SubElement(set_cTn, _pn('stCondLst'))
            set_condEl = etree.SubElement(set_stCond, _pn('cond'))
            set_condEl.set('delay', '0')

            set_tgtEl = etree.SubElement(set_cBhvr, _pn('tgtEl'))
            set_spTgt = etree.SubElement(set_tgtEl, _pn('spTgt'))
            set_spTgt.set('spid', spid)

            set_attrLst = etree.SubElement(set_cBhvr, _pn('attrNameLst'))
            set_attr = etree.SubElement(set_attrLst, _pn('attrName'))
            set_attr.text = 'style.visibility'

            set_to = etree.SubElement(setEl, _pn('to'))
            set_strVal = etree.SubElement(set_to, _pn('strVal'))
            set_strVal.set('val', 'visible')

    # Navigation conditions (prev / next)
    prevCond = etree.SubElement(seq, _pn('prevCondLst'))
    pc = etree.SubElement(prevCond, _pn('cond'))
    pc.set('evt', 'onPrev')
    pc.set('delay', '0')
    pc_tgt = etree.SubElement(pc, _pn('tgtEl'))
    etree.SubElement(pc_tgt, _pn('sldTgt'))

    nextCond = etree.SubElement(seq, _pn('nextCondLst'))
    nc = etree.SubElement(nextCond, _pn('cond'))
    nc.set('evt', 'onNext')
    nc.set('delay', '0')
    nc_tgt = etree.SubElement(nc, _pn('tgtEl'))
    etree.SubElement(nc_tgt, _pn('sldTgt'))


if __name__ == '__main__':
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Create 3 boxes
    boxes = []
    for i in range(3):
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(1 + i * 3), Inches(2),
            Inches(2), Inches(2),
        )
        shape.text = f'Box {i + 1}'
        boxes.append(shape)

    # Group: click 1 shows box 1, click 2 shows boxes 2 and 3 together
    click_groups = [
        [boxes[0].shape_id],
        [boxes[1].shape_id, boxes[2].shape_id],
    ]

    inject_appear_animations(slide._element, click_groups)

    prs.save('/tmp/test_animation.pptx')
    print('OK')
