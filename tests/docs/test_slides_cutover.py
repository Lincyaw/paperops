from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from paperops.slides import render_json

REPO_ROOT = Path(__file__).resolve().parents[2]


def test_legacy_slide_api_grep_is_empty():
    offenders: list[str] = []
    for relative in ["src", "examples"]:
        for path in (REPO_ROOT / relative).rglob("*"):
            if (
                path.suffix.lower() in {".pdf", ".pptx", ".png", ".jpg", ".jpeg"}
                or not path.is_file()
            ):
                continue
            text = path.read_text(encoding="utf-8", errors="ignore")
            for needle in ["Box(x=", "TextBlock(x=", "SlideBuilder"]:
                if needle in text:
                    offenders.append(f"{path}:{needle}")
    assert offenders == []


def test_mirrored_skills_match():
    for name in ["slidecraft", "slide-review", "talk-architect"]:
        claude = (REPO_ROOT / ".claude" / "skills" / name / "SKILL.md").read_text(
            encoding="utf-8"
        )
        codex = (REPO_ROOT / ".codex" / "skills" / name / "SKILL.md").read_text(
            encoding="utf-8"
        )
        assert claude == codex


def test_readme_and_quickstart_cover_python_json_and_mdx():
    readme = (REPO_ROOT / "README.md").read_text(encoding="utf-8")
    quickstart = (REPO_ROOT / "docs" / "quickstart-slides.md").read_text(
        encoding="utf-8"
    )

    assert "## Slides: 5-minute quickstart" in readme
    assert "### Python" in readme
    assert "### JSON" in readme
    assert "### MDX" in readme
    assert "LLM system prompt fragment" in quickstart
    assert "IR-first" in quickstart


def test_public_render_json_accepts_mdx_paths(tmp_path: Path):
    source = tmp_path / "quickstart.deck.mdx"
    source.write_text(
        """---
theme: minimal
sheet: keynote
---

# World-model RCA {.cover}

<Subtitle>Author once, render many variants</Subtitle>

---

## Why this API

<Grid cols="1fr 1fr" gap="md">
  <Text class="card">Use semantic components and classes.</Text>
  <Text class="card">Switch sheets without changing content.</Text>
</Grid>
""",
        encoding="utf-8",
    )
    out = tmp_path / "quickstart.pptx"

    render_json(source, out=out)

    assert len(Presentation(str(out)).slides) == 2


def test_changelog_records_breaking_slide_refactor():
    changelog = (REPO_ROOT / "CHANGELOG.md").read_text(encoding="utf-8")
    assert "breaking: slides API refactor" in changelog
