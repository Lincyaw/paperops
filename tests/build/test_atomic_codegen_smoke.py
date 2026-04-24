from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from paperops.slides.build import render_json


def test_render_all_atomic_components_via_single_deck(tmp_path: Path):
    payload = {
        "theme": "minimal",
        "sheet": "minimal",
        "slides": [
            {
                "type": "slide",
                "children": [
                    {"type": "text", "text": "text atom"},
                    {"type": "heading", "text": "heading atom"},
                    {"type": "title", "text": "title atom"},
                    {"type": "subtitle", "text": "subtitle atom"},
                    {"type": "box", "text": "box atom"},
                    {"type": "circle", "text": "◯"},
                    {"type": "line", "text": "line atom"},
                    {"type": "arrow", "text": "arrow atom"},
                    {"type": "divider", "props": {"orientation": "horizontal"}},
                    {"type": "badge", "text": "badge atom"},
                    {"type": "spacer", "props": {"size": "sm"}},
                    {"type": "icon", "props": {"name": "spark", "size": "sm"}},
                    {"type": "svg", "props": {"body": "<svg/>"}},
                    {"type": "image", "props": {"src": "placeholder.png"}},
                    {"type": "chart", "props": {"chart_type": "line"}},
                    {"type": "table", "props": {"headers": ["h1"], "rows": [["v1"]]}},
                    {"type": "roundedbox", "text": "rounded atom"},
                ]
            }
        ],
    }

    out = tmp_path / "all-atoms.pptx"
    render_json(payload, out=out)

    prs = Presentation(str(out))
    assert len(prs.slides) == 1
    slide = prs.slides[0]
    assert any(shape.has_table for shape in slide.shapes)
