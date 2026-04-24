from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from paperops.slides.build import render_json

pytest.importorskip("pptx")


SEMANTIC_COMPONENTS = [
    {"type": "card", "children": [{"type": "text", "text": "card body"}]},
    {"type": "kpi", "props": {"label": "DAU", "value": "125k", "delta": "+56%"}},
    {"type": "callout", "props": {"kind": "Insight", "text": "Keep signal clean"}},
    {"type": "quote", "props": {"text": "A good deck is clear.", "author": "Ops Team"}},
    {"type": "pullquote", "props": {"text": "Visibility drives trust.", "author": "Field notes"}},
    {"type": "keypoint", "props": {"number": "01", "title": "Start", "body": "Track evidence"}},
    {
        "type": "stepper",
        "props": {"steps": [{"label": "Collect"}, {"label": "Analyze"}]},
    },
    {"type": "timeline", "props": {"items": [{"date": "2026-01", "title": "Pilot"}]}},
    {"type": "figure", "props": {"chart_type": "bar", "caption": "Trend"}},
    {"type": "caption", "props": {"text": "A tiny caption"}},
    {"type": "spacer", "props": {"size": "md"}},
    {"type": "note", "children": ["speaker note"]},
]


def test_all_semantic_components_render_in_one_slide(tmp_path: Path) -> None:
    payload = {
        "theme": "minimal",
        "sheet": "minimal",
        "slides": [{"type": "slide", "children": SEMANTIC_COMPONENTS}],
    }
    output = tmp_path / "semantic-components.pptx"
    render_json(payload, out=output)

    assert output.exists() and output.stat().st_size > 0


def test_notes_are_written_to_speaker_notes_not_slide_shapes(tmp_path: Path) -> None:
    note_text = "Speaker note should never appear in slide body"
    payload = {
        "theme": "minimal",
        "sheet": "minimal",
        "slides": [{"type": "slide", "children": [{"type": "note", "children": [note_text]}]}],
    }
    output = tmp_path / "semantic-note.pptx"
    render_json(payload, out=output)

    prs = Presentation(str(output))
    assert len(prs.slides) == 1
    slide = prs.slides[0]
    assert note_text in slide.notes_slide.notes_text_frame.text
    assert all(note_text not in shape.text_frame.text for shape in slide.shapes if hasattr(shape, "text_frame"))
