from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from paperops.slides.build import (
    autofit_stage,
    codegen_stage,
    expand_stage,
    layout_stage,
    parse_stage,
    render_json,
    style_stage,
)
from paperops.slides.components.registry import ComponentError
from paperops.slides.ir.node import Node

MINIMAL_RAW_IR = {
    "theme": "minimal",
    "slides": [
        {
            "type": "slide",
            "class": "cover",
            "children": [
                {"type": "title", "text": "PaperOps Pipeline"},
                {"type": "subtitle", "text": "phase 2"},
            ],
        },
        {
            "type": "slide",
            "children": [
                {
                    "type": "flex",
                    "children": [
                        {"type": "box", "text": "A"},
                        {
                            "type": "kpi",
                            "props": {"label": "k", "value": "100", "delta": "+1%"},
                        },
                    ],
                },
                {
                    "type": "grid",
                    "style": {"cols": "1fr 1fr"},
                    "children": [
                        {"type": "text", "text": "Cell 1"},
                        {"type": "text", "text": "Cell 2"},
                    ],
                },
                {"type": "stack", "children": [{"type": "text", "text": "stack"}]},
                {
                    "type": "padding",
                    "style": {"padding": "sm"},
                    "children": [{"type": "text", "text": "padded"}],
                },
                {
                    "type": "layer",
                    "children": [{"type": "text", "text": "overlay"}],
                },
            ],
        },
    ],
}


def test_parse_stage_accepts_valid_json_ir():
    document = parse_stage(MINIMAL_RAW_IR)
    assert document.theme == "minimal"
    assert len(document.slides) == 2
    assert document.slides[0].type == "slide"


def test_style_stage_attaches_computed_styles_to_nodes():
    document = parse_stage(MINIMAL_RAW_IR)
    root, computed_styles = style_stage(document)
    assert root.type == "deck"
    assert root.children is not None and root.children[0].type == "slide"
    assert computed_styles
    first_slide = root.children[0]
    assert hasattr(first_slide, "computed_style")


def test_expand_stage_expands_semantic_components():
    document = parse_stage(MINIMAL_RAW_IR)
    expanded = expand_stage(document, strict=True)
    assert any(slide.type == "slide" for slide in expanded)
    source_types = {
        child.type
        for slide in expanded
        for child in (slide.children or [])
        if isinstance(child, Node)
    }
    assert "kpi" not in source_types


def test_render_stage_expands_kpi_into_cards():
    # Render-time style stage should include semantic expansion.
    deck = parse_stage(
        {
            "theme": "minimal",
            "slides": [
                {
                    "type": "slide",
                    "children": [
                        {"type": "kpi", "props": {"label": "DAU", "value": "125k"}},
                    ],
                }
            ],
        }
    )
    styled_root, _ = style_stage(deck, strict=True)
    first_slide = styled_root.children[0]
    assert isinstance(first_slide, object)
    slide_children = getattr(first_slide, "children", [])
    assert slide_children, "Expanded slide should still carry child nodes"
    child_types = {child.type for child in slide_children if hasattr(child, "type")}
    assert "kpi" not in child_types


def test_layout_stage_returns_slide_layout_roots():
    document = parse_stage(MINIMAL_RAW_IR)
    root, _ = style_stage(document)
    slide_layouts, issues = layout_stage(root, theme=document.theme)
    assert len(slide_layouts) == 2
    assert issues == []
    assert all(source_slide.type == "slide" for source_slide, _ in slide_layouts)
    assert all(hasattr(layout_node, "_region") for _, layout_node in slide_layouts)


def test_autofit_stage_is_stage_boundary_compatible():
    document = parse_stage(MINIMAL_RAW_IR)
    root, _ = style_stage(document)
    slide_layouts, _ = layout_stage(root, theme=document.theme)
    fixed = autofit_stage(slide_layouts, theme=document.theme)
    assert len(fixed) == len(slide_layouts)
    assert all(source.type == "slide" for source, _ in fixed)


def test_codegen_stage_writes_openable_pptx(tmp_path: Path):
    document = parse_stage(MINIMAL_RAW_IR)
    root, _ = style_stage(document)
    slide_layouts, _ = layout_stage(root, theme=document.theme)
    out = tmp_path / "phase2-smoke.pptx"
    codegen_stage(autofit_stage(slide_layouts), theme=document.theme, out=out)
    assert out.exists()
    assert Presentation(str(out)) is not None


def test_render_json_stage_wires_the_full_pipeline(tmp_path: Path):
    out = tmp_path / "pipeline-full.pptx"
    returned = render_json(MINIMAL_RAW_IR, out=out)
    assert returned == out
    prs = Presentation(str(out))
    assert len(prs.slides) == 2


def test_render_json_enforces_semantic_prop_requirements(tmp_path: Path):
    with pytest.raises(ComponentError, match="MISSING_REQUIRED_PROP"):
        render_json(
            {
                "theme": "minimal",
                "slides": [
                    {
                        "type": "slide",
                        "children": [{"type": "kpi", "props": {"label": "DAU"}}],
                    }
                ],
            },
            out=tmp_path / "missing-prop.pptx",
        )

    with pytest.raises(ComponentError, match="UNKNOWN_PROP"):
        render_json(
            {
                "theme": "minimal",
                "slides": [
                    {
                        "type": "slide",
                        "children": [
                            {
                                "type": "kpi",
                                "props": {
                                    "label": "DAU",
                                    "value": "125k",
                                    "delta": "+56%",
                                    "unknown": "x",
                                },
                            }
                        ],
                    }
                ],
            },
            out=tmp_path / "unknown-prop.pptx",
        )
