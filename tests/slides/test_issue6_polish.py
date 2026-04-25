from __future__ import annotations

import json
from pathlib import Path
from zipfile import ZipFile

import pytest
from lxml import etree
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE

from paperops.slides.build import (
    layout_stage,
    parse_stage,
    render_diagnostics,
    render_json,
    style_stage,
)
from paperops.slides.core.theme import themes
from paperops.slides.ir.validator import StructuredValidationError, validate_ir_document

PPT_NS = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}


def _iter_layout_nodes(node):
    yield node
    if hasattr(node, "children") and node.children:
        for child in node.children:
            if child is not None:
                yield from _iter_layout_nodes(child)
    if hasattr(node, "child") and node.child is not None:
        yield from _iter_layout_nodes(node.child)
    if hasattr(node, "iter_items"):
        for item in node.iter_items():
            if item is not None and item.child is not None:
                yield from _iter_layout_nodes(item.child)


def _shape_with_text(slide, text: str):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and text in shape.text:
            return shape
    raise AssertionError(f"No shape contains text {text!r}")


def _slide_xml_texts(path: Path, member: str) -> etree._Element:
    with ZipFile(path) as archive:
        return etree.fromstring(archive.read(member))


def test_shrink_overflow_sets_ppt_autofit(tmp_path: Path):
    out = tmp_path / "shrink.pptx"
    render_json(
        {
            "theme": "minimal",
            "slides": [
                {
                    "type": "slide",
                    "children": [
                        {
                            "type": "text",
                            "text": "This is a deliberately long string that should shrink to fit.",
                            "style": {
                                "width": 2.2,
                                "height": 0.35,
                                "overflow": "shrink",
                            },
                        }
                    ],
                }
            ],
        },
        out=out,
    )

    shape = _shape_with_text(Presentation(str(out)).slides[0], "deliberately long")
    assert shape.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def test_reflow_overflow_creates_continuation_slide(tmp_path: Path):
    out = tmp_path / "reflow.pptx"
    render_json(
        {
            "theme": "minimal",
            "meta": {"lang": "en-US"},
            "slides": [
                {
                    "type": "slide",
                    "children": [
                        {"type": "title", "text": "Overflow Demo"},
                        {
                            "type": "absolute",
                            "children": [
                                {
                                    "type": "prose",
                                    "style": {
                                        "left": 0.8,
                                        "top": 1.8,
                                        "width": 3.4,
                                        "height": 0.7,
                                        "overflow": "reflow",
                                    },
                                    "children": [
                                        "Paragraph one. " * 2,
                                        "Paragraph two. " * 2,
                                        "Paragraph three. " * 2,
                                        "Paragraph four. " * 2,
                                    ],
                                }
                            ],
                        },
                    ],
                }
            ],
        },
        out=out,
    )

    prs = Presentation(str(out))
    assert len(prs.slides) >= 2
    assert _shape_with_text(prs.slides[0], "Overflow Demo")
    assert _shape_with_text(prs.slides[1], "Overflow Demo (cont.)")


def test_clip_overflow_disables_autofit(tmp_path: Path):
    out = tmp_path / "clip.pptx"
    render_json(
        {
            "theme": "minimal",
            "slides": [
                {
                    "type": "slide",
                    "children": [
                        {
                            "type": "code",
                            "text": "very long code block " * 10,
                            "style": {"width": 2.0, "height": 0.3, "overflow": "clip"},
                        }
                    ],
                }
            ],
        },
        out=out,
    )

    shape = _shape_with_text(Presentation(str(out)).slides[0], "very long code")
    assert shape.text_frame.auto_size == MSO_AUTO_SIZE.NONE


def test_error_overflow_raises_structured_error(tmp_path: Path):
    with pytest.raises(StructuredValidationError) as excinfo:
        render_json(
            {
                "theme": "minimal",
                "slides": [
                    {
                        "type": "slide",
                        "children": [
                            {
                                "type": "absolute",
                                "children": [
                                    {
                                        "type": "text",
                                        "text": "cannot fit " * 30,
                                        "style": {
                                            "left": 0.8,
                                            "top": 1.5,
                                            "width": 2.0,
                                            "height": 0.2,
                                            "overflow": "error",
                                        },
                                    }
                                ],
                            }
                        ],
                    }
                ],
            },
            out=tmp_path / "error.pptx",
        )

    payload = excinfo.value.to_dict()
    assert payload["errors"][0]["code"] == "OVERFLOW_UNRECOVERABLE"
    assert (
        payload["errors"][0]["needed_height"] > payload["errors"][0]["available_height"]
    )


def test_baseline_snap_and_align_to_sibling():
    document = parse_stage(
        {
            "theme": "minimal",
            "slides": [
                {
                    "type": "slide",
                    "style": {"baseline-snap": True},
                    "children": [
                        {
                            "type": "absolute",
                            "children": [
                                {
                                    "type": "title",
                                    "class": "title",
                                    "text": "Baseline",
                                    "style": {
                                        "left": 0.3,
                                        "top": 0.23,
                                        "width": 2.0,
                                        "height": 0.31,
                                    },
                                },
                                {
                                    "type": "text",
                                    "text": "Aligned body",
                                    "style": {
                                        "left": 0.3,
                                        "top": 1.07,
                                        "width": 2.0,
                                        "height": 0.31,
                                        "align-to": "sibling.title:bottom",
                                    },
                                },
                            ],
                        }
                    ],
                }
            ],
        }
    )
    styled_root, _ = style_stage(document)
    slide_layouts, _ = layout_stage(styled_root, theme=document.theme)
    _, layout_root = slide_layouts[0]

    title_region = None
    body_region = None
    for layout_node in _iter_layout_nodes(layout_root):
        source = getattr(layout_node, "_ir_node", None)
        if source is None or getattr(layout_node, "_region", None) is None:
            continue
        if source.type == "title":
            title_region = layout_node._region
        elif source.type == "text":
            body_region = layout_node._region

    assert title_region is not None and body_region is not None
    baseline = themes.minimal.baseline
    assert round(title_region.top / baseline) * baseline == pytest.approx(
        title_region.top
    )
    assert round(body_region.top / baseline) * baseline == pytest.approx(
        body_region.top
    )
    assert body_region.top == pytest.approx(
        round((title_region.top + title_region.height) / baseline) * baseline
    )


def test_animation_timing_and_notes_are_written(tmp_path: Path):
    out = tmp_path / "animation.pptx"
    render_json(
        {
            "theme": "minimal",
            "slides": [
                {
                    "type": "slide",
                    "children": [
                        {
                            "type": "note",
                            "children": ["Remember to mention stagger timing."],
                        },
                        {
                            "type": "text",
                            "text": "First",
                            "style": {
                                "animate": "fade-up",
                                "animate-group": "intro",
                                "animate-trigger": "on-click",
                                "stagger": "fast",
                            },
                        },
                        {
                            "type": "text",
                            "text": "Second",
                            "style": {
                                "animate": "fade-up",
                                "animate-group": "intro",
                                "animate-trigger": "on-click",
                                "stagger": "fast",
                            },
                        },
                        {
                            "type": "text",
                            "text": "Third",
                            "style": {
                                "animate": "zoom-in",
                                "animate-group": "outro",
                                "animate-trigger": "after-previous",
                            },
                        },
                    ],
                }
            ],
        },
        out=out,
    )

    slide_xml = _slide_xml_texts(out, "ppt/slides/slide1.xml")
    notes_xml = _slide_xml_texts(out, "ppt/notesSlides/notesSlide1.xml")
    delays = slide_xml.xpath(".//p:timing//p:cond/@delay", namespaces=PPT_NS)
    effects = slide_xml.xpath(".//p:timing//p:animEffect", namespaces=PPT_NS)
    notes_text = "".join(
        notes_xml.xpath(
            ".//a:t/text()",
            namespaces={
                **PPT_NS,
                "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            },
        )
    )

    assert effects
    assert "indefinite" in delays
    assert str(int(themes.minimal.duration["fast"] * 1000)) in delays
    assert "Remember to mention stagger timing." in notes_text


def test_structured_validator_matches_fixture_and_supports_text_json_output():
    fixture_dir = Path("tests/fixtures/issue6")
    payload = json.loads((fixture_dir / "bad_ir.json").read_text(encoding="utf-8"))
    expected = json.loads(
        (fixture_dir / "bad_ir.expected.json").read_text(encoding="utf-8")
    )

    report = validate_ir_document(payload, strict=True)
    assert report.to_dict() == expected
    assert json.loads(report.to_json()) == expected
    pretty = report.to_pretty_text()
    assert "UNKNOWN_STYLE_KEY" in pretty
    assert "UNRESOLVED_MACRO_VAR" in pretty


def test_trace_option_returns_matched_rules_and_computed_keys():
    diagnostics = render_diagnostics(
        {
            "theme": "minimal",
            "styles": {".hero": {"animate": "fade-up", "duration": "fast"}},
            "slides": [
                {
                    "type": "slide",
                    "children": [
                        {"type": "text", "class": "hero", "text": "Trace me"},
                    ],
                }
            ],
        },
        trace=True,
    )

    assert diagnostics["trace"]
    assert any(".hero" in entry["matched_rules"] for entry in diagnostics["trace"])
    assert any(
        "animate" in entry["computed_style_keys"] for entry in diagnostics["trace"]
    )
