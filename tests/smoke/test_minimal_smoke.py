from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation

from paperops.slides import build as build_module


def test_smoke_render_minimal_deck_from_json(tmp_path: Path):
    fixture = Path(__file__).with_name("minimal_deck.json")
    deck_payload = json.loads(fixture.read_text(encoding="utf-8"))
    output = tmp_path / "minimal_smoke.pptx"

    rendered = build_module.render_json(deck_payload, out=output)
    assert rendered == output
    assert output.exists()

    prs = Presentation(str(output))
    assert len(prs.slides) == 3
    assert any(slide.shapes for slide in prs.slides)
