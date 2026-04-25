from __future__ import annotations

import shutil
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from pptx import Presentation

from examples.self_intro.build import (
    IMPLEMENTED_SLIDE_COUNT,
    OUTPUT_FILE,
    SLIDE_TITLES,
    build_presentation,
    make_theme,
)


def _slide_texts(prs: Presentation) -> list[str]:
    texts: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
        texts.append("\n".join(parts))
    return texts


def test_self_intro_exports_ir_first_metadata():
    assert IMPLEMENTED_SLIDE_COUNT == len(SLIDE_TITLES) == 5
    assert SLIDE_TITLES[0] == "Diagnostic intelligence should be trained around world models"
    assert make_theme() == "minimal"


def test_build_self_intro_deck(tmp_path: Path):
    out_path = tmp_path / OUTPUT_FILE.name
    prs = build_presentation(output_path=out_path, render_preview=False)

    assert out_path.exists()
    assert len(prs.slides) == len(SLIDE_TITLES)

    slide_texts = _slide_texts(prs)
    assert "Trusted diagnosis needs evidence" in slide_texts[-1]
    assert "9,152" in "\n".join(slide_texts)


def test_build_self_intro_preview_when_requested(tmp_path: Path):
    if shutil.which("soffice") is None or shutil.which("pdftoppm") is None:
        return

    out_path = tmp_path / OUTPUT_FILE.name
    build_presentation(output_path=out_path, render_preview=True)

    preview_dir = tmp_path / "preview"
    preview_files = sorted(preview_dir.glob("slide_*.png"))
    assert len(preview_files) == len(SLIDE_TITLES)
