from __future__ import annotations

import shutil
import subprocess
import sys
from pathlib import Path

from pptx import Presentation

from paperops.slides.preview import render_slide_preview_powerpoint

REPO_ROOT = Path(__file__).resolve().parents[2]
EXAMPLES = [
    (
        REPO_ROOT / "examples" / "4.15-talk" / "presentation.py",
        "talk_4_15_ir_first.pptx",
    ),
    (
        REPO_ROOT / "examples" / "4.15-talk" / "pptx_restructure.py",
        "variants/talk_4_15_seminar.pptx",
    ),
    (REPO_ROOT / "examples" / "fse26" / "presentation.py", "fse26_ir_first.pptx"),
    (
        REPO_ROOT / "examples" / "communication_friction" / "build.py",
        "communication_friction_ir_first.pptx",
    ),
    (
        REPO_ROOT / "examples" / "self_intro" / "build.py",
        "diagnostic_world_model_ir_first.pptx",
    ),
]


def test_ir_first_examples_build_pptx():
    for script_path, expected_name in EXAMPLES:
        subprocess.run([sys.executable, str(script_path)], check=True, cwd=REPO_ROOT)
        output_path = script_path.parent / expected_name
        assert output_path.exists(), f"missing output for {script_path}"
        prs = Presentation(str(output_path))
        assert len(prs.slides) >= 2


def test_hku_talk_opens_in_libreoffice_preview_pipeline(tmp_path: Path):
    if shutil.which("soffice") is None or shutil.which("pdftoppm") is None:
        return

    script = REPO_ROOT / "examples" / "4.15-talk" / "presentation.py"
    subprocess.run([sys.executable, str(script)], check=True, cwd=REPO_ROOT)
    pptx_path = script.parent / "talk_4_15_ir_first.pptx"

    preview_paths = render_slide_preview_powerpoint(
        str(pptx_path), str(tmp_path / "preview"), dpi=72
    )

    assert len(preview_paths) == 5
    assert all(Path(path).stat().st_size > 0 for path in preview_paths)


def test_gallery_script_renders_six_distinct_variants(tmp_path: Path):
    script = REPO_ROOT / "examples" / "gallery" / "render_gallery_variants.py"
    out_dir = tmp_path / "gallery"
    subprocess.run(
        [sys.executable, str(script), "--output-dir", str(out_dir)],
        check=True,
        cwd=REPO_ROOT,
    )
    outputs = sorted(out_dir.glob("gallery_*.pptx"))
    assert len(outputs) == 6
    sizes = {path.stat().st_size for path in outputs}
    assert len(sizes) == len(outputs)
