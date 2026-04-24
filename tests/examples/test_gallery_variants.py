from __future__ import annotations

import hashlib
import json
from pathlib import Path

from pptx import Presentation

from paperops.slides.build import render_json


GALLERY_CONTENT = Path(__file__).resolve().parents[2] / "examples" / "gallery" / "gallery_content.json"
SHEETS = ["minimal", "academic", "seminar", "keynote", "whitepaper", "pitch"]


def _read_content() -> dict:
    return json.loads(GALLERY_CONTENT.read_text(encoding="utf-8"))


def _checksum(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def test_render_all_sheet_variants_from_gallery_content(tmp_path: Path):
    base = _read_content()
    outputs: list[Path] = []

    for sheet in SHEETS:
        payload = dict(base)
        payload["sheet"] = sheet
        out = tmp_path / f"gallery-{sheet}.pptx"
        render_json(payload, out=out)
        outputs.append(out)

        prs = Presentation(str(out))
        assert len(prs.slides) == len(base["slides"])

    hashes = {_checksum(out) for out in outputs}
    assert len(hashes) == len(outputs), "variant outputs should differ by sheet"
